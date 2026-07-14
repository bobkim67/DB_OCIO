#!/usr/bin/env python3
"""P0' 파일럿 — s7(성과 리뷰) 슬라이드: A4 가로(29.7x21cm) + matplotlib + python-pptx(Route B).

레이아웃 (2026-07-13 사용자 지시):
  - 슬라이드 = A4 가로 (297x210mm). 원본 16:9(1600x900) base 를 흰 행 삽입으로
    1600x1131px 재구성(헤더/푸터 무왜곡), 콘텐츠는 헤더~푸터라인 사이 선형 재배치.
  - 본문 텍스트(불릿·표·타이틀바) = 10pt 고정. 제목패널(성과 리뷰·기준일)만 비례 크기.
편집 가능: 제목·기준일·불릿·타이틀바(도형)·표 2종(네이티브). 이미지 = 성과추이 차트만.
한글 폰트: a:latin + a:ea + a:cs 모두 Pretendard (테마 맑은고딕 폴백 방지).

입력: reporting/reference/fund_07G07_YTD.json / _2Q.json (레포 스냅샷 — 원본 대조 목적)
출력: reporting/pilot/out/pilot_s7.pptx
"""
import json
import datetime
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
from matplotlib import pyplot as plt, font_manager

HERE = Path(__file__).parent            # reporting/pilot
ROOT = HERE.parent                      # reporting
OUT = HERE / 'out'
FONTS = ROOT / 'template' / 'fonts'
BASE = ROOT / 'template' / 'base' / 'base_slide07.png'

for _f in FONTS.glob('*.otf'):
    font_manager.fontManager.addfont(str(_f))
plt.rcParams['font.family'] = 'Pretendard'

# ── A4 가로 좌표계 ──────────────────────────────────────────────
# 가로 1600px = 29.7cm → EMU/px = 10692000/1600 = 6682.5, 세로 21cm = 1131.3px
EMU_A4_W, EMU_A4_H = 10_692_000, 7_560_000
EMU_PER_PX = EMU_A4_W / 1600            # 6682.5
PX_H = 1131                             # base A4 세로 px
PT_PER_PX = 72 / (1600 / (EMU_A4_W / 914_400))   # 0.5262 (제목패널 비례 환산용)
BODY_PT = 12                            # 본문 고정 폰트 (2026-07-13 사용자 지시, 10→12pt)

CANVAS_OFF = (40, 192)                  # 원본 compose(): s7 캔버스 슬라이드 좌표 (x 그대로)

# 세로 재배치: 헤더끝(196)~푸터라인(원본824→A4 1055) 구간 선형 스트레치
_FOOT_OLD, _FOOT_NEW, _HDR = 824, 824 + (PX_H - 900), 196     # 1055
SV = (_FOOT_NEW - _HDR) / (_FOOT_OLD - _HDR)                  # 1.3678


def remap(y):
    """900px 레이아웃의 y → A4(1131px) 콘텐츠존 y."""
    return round(_HDR + (y - _HDR) * SV) if y >= _HDR else y


HDR = '5B9BD5'; Z1 = 'D2DEEF'; Z2 = 'EAEFF7'
C_FUND = '#2E5E9E'; C_BM = '#E0A800'; C_EXC = '#B9C9E8'


# ──────────────────────────── 데이터 (원본 fund_calcs 포트) ────────────────────────────
def fund_calcs():
    fund = json.loads((ROOT / 'reference' / 'fund_07G07_YTD.json').read_text(encoding='utf-8'))
    fund2q = json.loads((ROOT / 'reference' / 'fund_07G07_2Q.json').read_text(encoding='utf-8'))
    sm, sm2 = fund['summary'], fund2q['summary']
    ser = fund['series']
    end = ser[-1]
    end_d = datetime.date(*map(int, end['date'].split('-')))
    py, pm = (end_d.year, end_d.month - 1) if end_d.month > 1 else (end_d.year - 1, 12)
    last_dom = (datetime.date(py, pm, 28) + datetime.timedelta(days=4)).replace(day=1) - datetime.timedelta(days=1)
    m1_target = datetime.date(py, pm, min(end_d.day, last_dom.day))
    m1 = next(r for r in reversed(ser) if r['date'] <= m1_target.isoformat())
    return {
        'asof': sm['end'],
        'ytd_f': sm['fund_return_pct'], 'ytd_b': sm['bm_return_pct'], 'ytd_x': sm['excess_pct'],
        'q_f': sm2['fund_return_pct'], 'q_b': sm2['bm_return_pct'], 'q_x': sm2['excess_pct'],
        'm1_f': round((end['nav'] / m1['nav'] - 1) * 100, 2),
        'm1_b': round((end['bm'] / m1['bm'] - 1) * 100, 2),
        'si_f': round((end['nav'] / 1000 - 1) * 100, 2),
        'si_b': round((end['bm'] / 1000 - 1) * 100, 2),
        'series': ser,
    }


# ──────────────────────────── base A4 재구성 ────────────────────────────
def clean_base():
    """1600x900 base → 1600x1131 A4: 제목 제거 + 밴드 균일화 + 흰 행 삽입(카드 구간)."""
    from PIL import Image, ImageDraw
    import numpy as np
    im = Image.open(BASE).convert('RGB')
    d = ImageDraw.Draw(im)
    d.rectangle([25, 12, 1150, 118], fill='white')   # "성과 리뷰" 영역 (편집 텍스트로 대체)
    # 우상단 저해상 로고 제거 — 고해상 ki-logo.png(2467x607) 를 pptx 그림 개체로 별도 삽입
    # (base 에 굽지 않음 → 선명 + PPT 에서 조정 가능. 2026-07-13 사용자 지시)
    d.rectangle((1460, 12, 1576, 37), fill='white')
    a = np.array(im)
    # 음영 밴드 전폭 균일화 + 코너 딥(y184~192) 평탄화 — 클린 열 x1300 프로파일 복제
    a[118:196, :] = a[118:196, 1300:1301]
    # A4 세로 확장: y500(카드 내 균일 백색 행, 사전 검증) 에 231행 삽입 → 헤더/푸터 무왜곡
    filler = np.repeat(a[500:501, :, :], PX_H - 900, axis=0)
    a = np.vstack([a[:500], filler, a[500:]])
    assert a.shape[0] == PX_H
    OUT.mkdir(parents=True, exist_ok=True)
    p = OUT / 'base_slide07_a4.png'
    Image.fromarray(a).save(p)
    return p


# ──────────────────────────── 차트: 연초 이후 성과 추이 (세로 스트레치) ────────────────────────────
def gen_perf_chart(fc):
    # 배치: y = remap(482)=587, 하단 여유 → 푸터라인(1055) 침범 금지
    IMG_X0 = 48                          # 캔버스 x (슬라이드 x=88)
    W = 724
    H = round(340 * SV)                  # 465
    PL, PT = 48, 10
    CW, CH = 640, round(280 * SV)        # 383
    YMAX, YMIN = 14, -2                  # 펀드 최대 +13.01% 상단 잘림 방지
    ser = fc['series']
    n0, b0 = ser[0]['nav'], ser[0]['bm']
    N = len(ser)

    def X(i): return PL + i / (N - 1) * CW
    def Y(v): return PT + (YMAX - v) / (YMAX - YMIN) * CH

    fig = plt.figure(figsize=(W / 72, H / 72), dpi=144)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W); ax.set_ylim(H, 0); ax.axis('off')
    for v in range(YMIN, YMAX + 1, 2):
        ax.plot([PL, PL + CW], [Y(v)] * 2, color='#999999' if v == 0 else '#E8E8E8', lw=1, zorder=1)
        ax.text(PL - 6, Y(v), f'{v}%', ha='right', va='center', fontsize=13, color='#666')
    xs = [X(i) for i in range(N)]
    fv = [(r['nav'] / n0 - 1) * 100 for r in ser]
    bv = [(r['bm'] / b0 - 1) * 100 for r in ser]
    ev = [f - b for f, b in zip(fv, bv)]
    ax.fill_between(xs, [Y(0)] * N, [Y(e) for e in ev], color=C_EXC, alpha=0.8, lw=0, zorder=2)
    ax.plot(xs, [Y(v) for v in fv], color=C_FUND, lw=2, zorder=3)
    ax.plot(xs, [Y(v) for v in bv], color=C_BM, lw=2, zorder=3)
    for i, r in enumerate(ser):
        if r['date'][8:10] == '01':
            ax.text(X(i) - 14, PT + CH + 6, f"{int(r['date'][5:7])}월",
                    ha='left', va='top', fontsize=13, color='#666')
    ly = PT + CH + 42
    lx = 158
    ax.add_patch(plt.Rectangle((lx, ly - 5.5), 22, 11, facecolor=C_EXC, edgecolor='none'))
    ax.text(lx + 27, ly, '초과성과(%p)', ha='left', va='center', fontsize=13.5, color='#555')
    lx += 27 + 90 + 22
    ax.plot([lx, lx + 22], [ly] * 2, color=C_FUND, lw=3)
    ax.text(lx + 27, ly, '펀드', ha='left', va='center', fontsize=13.5, color='#555')
    lx += 27 + 32 + 22
    ax.plot([lx, lx + 22], [ly] * 2, color=C_BM, lw=3)
    ax.text(lx + 27, ly, 'BM', ha='left', va='center', fontsize=13.5, color='#555')
    p = OUT / 's7_perf_chart.png'
    fig.savefig(p, facecolor='white')
    plt.close(fig)
    return p, (IMG_X0, W, H)


# ──────────────────────────── pptx 조립 (Route B, A4) ────────────────────────────
def build_pptx(fc, chart_png, chart_geom):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.lang import MSO_LANGUAGE_ID
    from pptx.oxml.ns import qn

    def E(px): return Emu(round(px * EMU_PER_PX))

    def set_ko_font(font, family):
        """a:latin 만으로는 한글이 테마 EA(맑은고딕)로 폴백 — a:ea/a:cs 도 지정."""
        font.language_id = MSO_LANGUAGE_ID.KOREAN
        rPr = font._rPr
        latin = rPr.find(qn('a:latin')) if rPr is not None else None
        if latin is None:
            return
        for tag in ('a:cs', 'a:ea'):        # addnext 역순 삽입 → latin, ea, cs 순서
            e = rPr.makeelement(qn(tag), {'typeface': family})
            latin.addnext(e)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(EMU_A4_W), Emu(EMU_A4_H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])     # blank

    sl.shapes.add_picture(str(clean_base()), 0, 0, Emu(EMU_A4_W), Emu(EMU_A4_H))
    OX, OY = CANVAS_OFF

    # 우상단 로고 — 고해상 원본(web/public/ki-logo.png, 2467x607) 별도 그림 개체
    ki_logo = ROOT.parent / 'web' / 'public' / 'ki-logo.png'
    LOGO_W = 150
    LOGO_H = round(LOGO_W * 607 / 2467)               # 종횡비 고정 (37px)
    sl.shapes.add_picture(str(ki_logo), E(1600 - 27 - LOGO_W), E(14), E(LOGO_W), E(LOGO_H))

    def add_text(px_x, px_y, px_w, px_h, text, pt_size, color, bold=False,
                 family='Pretendard', align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = sl.shapes.add_textbox(E(px_x), E(px_y), E(px_w), E(px_h))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.name = family; r.font.size = Pt(pt_size)
        r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
        set_ko_font(r.font, family)
        return tb

    # ── 제목패널 (회색 음영 헤더 — 비례 크기 유지, 10pt 규칙 제외) ──
    add_text(36, 22, 1000, 96, '성과 리뷰', 76 * PT_PER_PX, '000000', family='Pretendard Black')
    y, m, d = fc['asof'].split('-')
    add_text(55, 133, 700, 34, f'기준일: {y}년 {int(m)}월 {int(d)}일', 26 * PT_PER_PX,
             '7B401F', bold=True)

    # ── 본문 (전부 10pt) ──
    b1 = (f"· 연초 이후 수익률 {fc['ytd_f']:+.2f}%로 BM({fc['ytd_b']:+.2f}%) "
          f"대비 {fc['ytd_x']:+.2f}%p 초과 성과")
    b2 = (f"· 2분기 수익률: 펀드 {fc['q_f']:+.2f}%, BM {fc['q_b']:+.2f}% "
          f"(초과성과 {fc['q_x']:+.2f}%p)")
    # 불릿 x = 표1 시작점(OX+52)과 좌측 정렬 (2026-07-13 사용자 지시)
    add_text(OX + 52, remap(OY + 18), 1200, 30, b1, BODY_PT, '222222')
    add_text(OX + 52, remap(OY + 50), 1200, 30, b2, BODY_PT, '222222')

    PBAR_H = round(34 * SV)                              # 47
    def add_pbar(px_x, px_y, px_w, text):
        sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(px_x), E(px_y), E(px_w), E(PBAR_H))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(HDR)
        sh.line.fill.background()
        sh.shadow.inherit = False
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.name = 'Pretendard'; r.font.size = Pt(BODY_PT)
        r.font.bold = True; r.font.color.rgb = RGBColor.from_string('FFFFFF')
        set_ko_font(r.font, 'Pretendard')

    add_pbar(OX + 48, remap(OY + 246), 724, '연초 이후 성과 추이')
    add_pbar(OX + 788, remap(OY + 246), 694, '기여수익률 분석(단위: %, %p)')

    # ── 네이티브 표 (본문 10pt) ──
    def add_table(px_x, px_y, col_w_px, row_h_px, rows_spec):
        from lxml import etree
        n_r, n_c = len(rows_spec), len(col_w_px)
        gf = sl.shapes.add_table(n_r, n_c, E(px_x), E(px_y),
                                 E(sum(col_w_px)), E(sum(row_h_px)))
        tbl = gf.table
        tblPr = tbl._tbl.tblPr
        tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'   # No Style, No Grid
        for i, w in enumerate(col_w_px):
            tbl.columns[i].width = E(w)
        for i, h in enumerate(row_h_px):
            tbl.rows[i].height = E(h)
        for ri, (fill, cells) in enumerate(rows_spec):
            for ci, (text, bold, color) in enumerate(cells):
                cell = tbl.cell(ri, ci)
                if fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(fill)
                if ri < n_r - 1:                          # 행 사이 2px 흰 경계
                    tcPr = cell._tc.get_or_add_tcPr()
                    ln = etree.SubElement(tcPr, qn('a:lnB'))
                    ln.set('w', '15240'); ln.set('cap', 'flat')
                    sf = etree.SubElement(ln, qn('a:solidFill'))
                    etree.SubElement(sf, qn('a:srgbClr')).set('val', 'FFFFFF')
                    tcPr.insert(0, ln)
                # 셀 세로 가운데: 표 셀은 bodyPr 아닌 tcPr@anchor 가 적용됨 (2026-07-13 fix)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = cell.text_frame
                tf.margin_left = tf.margin_right = Emu(round(2 * EMU_PER_PX))
                tf.margin_top = tf.margin_bottom = 0
                tf.word_wrap = False
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                p.font.name = 'Pretendard'; p.font.size = Pt(BODY_PT)   # 빈 셀 행높이 부풀림 방지
                set_ko_font(p.font, 'Pretendard')
                if text:
                    r = p.add_run(); r.text = text
                    r.font.name = 'Pretendard'; r.font.size = Pt(BODY_PT)
                    r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
                    set_ko_font(r.font, 'Pretendard')
        return gf

    W = 'FFFFFF'; K = '222222'
    TH, TD = round(38 * SV), round(36 * SV)              # 52 / 49

    # 표1: 기간수익률
    add_table(
        OX + 52, remap(OY + 96), [230, 240, 240, 240, 240, 240], [TH, TD, TD],
        [
            (HDR, [(t, True, W) for t in ['기간', '1개월', '3개월', '6개월', '연초 이후', '설정 이후']]),
            (Z1, [('펀드(%)', True, K)] + [(f'{v:.2f}', False, K)
                  for v in [fc['m1_f'], fc['q_f'], fc['ytd_f'], fc['ytd_f'], fc['si_f']]]),
            (Z2, [('BM(%)', True, K)] + [(f'{v:.2f}', False, K)
                  for v in [fc['m1_b'], fc['q_b'], fc['ytd_b'], fc['ytd_b'], fc['si_b']]]),
        ])

    # 표2: 기여수익률
    classes = ['국내주식', '해외주식', '국내채권', '해외채권', '대체투자', '기타']
    rows = [(HDR, [(t, True, W) for t in ['연초 이후', '펀드', 'BM', 'Active']])]
    for i, name in enumerate(classes):
        rows.append((Z1 if i % 2 == 0 else Z2,
                     [(name, False, K), ('', False, K), ('', False, K), ('', False, K)]))
    rows.append((Z1, [('합계', True, K)] + [(f'{v:+.2f}', True, K)
                 for v in [fc['ytd_f'], fc['ytd_b'], fc['ytd_x']]]))
    add_table(OX + 788, remap(OY + 288), [190, 168, 168, 168], [TH] + [TD] * 7, rows)

    # 차트 이미지
    cx0, cw, ch = chart_geom
    sl.shapes.add_picture(str(chart_png), E(OX + cx0), E(remap(OY + 290)), E(cw), E(ch))

    out = OUT / 'pilot_s7.pptx'
    try:
        prs.save(str(out))
    except PermissionError:                      # PowerPoint 에 열려있으면 버전 폴백
        out = OUT / 'pilot_s7_v2.pptx'
        prs.save(str(out))
    return out


def main():
    fc = fund_calcs()
    chart_png, chart_geom = gen_perf_chart(fc)
    out = build_pptx(fc, chart_png, chart_geom)
    print(f'OK: {out}')
    print(f'  A4: {EMU_A4_W / 360000:.1f}x{EMU_A4_H / 360000:.1f}cm | SV={SV:.4f} | '
          f'footer line y={_FOOT_NEW}')


if __name__ == '__main__':
    main()
