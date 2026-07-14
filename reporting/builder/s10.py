"""슬라이드 10 — MSCI KR: 한국 가격지수 vs 12M 선행 PER (이중축, 이미지)."""
import datetime
import math

from .common import (OUT, E, EX, add_text, slide_scaffold, plt, BODY_PT,
                     PT_PER_PX, PP_ALIGN, kdate)

C_PX, C_PE = '#C0392B', '#2E6DB4'


def _dnum(d):
    y, m, dd = map(int, d.split('-'))
    return datetime.date(y, m, dd).toordinal()


def gen_chart(kr, end):
    W, H = 1520, 770
    PL, PT_, PW, PH = 78, 40, 1372, 600
    ser = kr['series']            # (date, pe, eps)
    px = kr['px']                 # (date, v)
    T0, T1 = _dnum(ser[0][0]), _dnum(ser[-1][0])

    pxv = [v for _, v in px]
    pev = [p for _, p, _ in ser]

    # 나이스 스텝 8분할 (좌 437.5 같은 지저분한 눈금 방지)
    def nice8(vmin, vmax, steps):
        for st in steps:
            bot = math.floor(vmin / st) * st
            if bot + st * 8 >= vmax:
                return bot, bot + st * 8, st
        st = steps[-1]
        bot = math.floor(vmin / st) * st
        return bot, math.ceil(vmax / st) * st, st

    L_BOT, L_TOP, _ = nice8(min(pxv), max(pxv), (250, 500, 1000, 2000))
    L_BOT = max(0, L_BOT)
    R_BOT, R_TOP, _ = nice8(min(pev), max(pev), (1, 2, 2.5, 5))
    R_BOT = max(0, R_BOT)
    n_div = 8

    def X(d): return PL + (_dnum(d) - T0) / (T1 - T0) * PW
    def YL(v): return PT_ + (L_TOP - v) / (L_TOP - L_BOT) * PH
    def YR(p): return PT_ + (R_TOP - p) / (R_TOP - R_BOT) * PH

    fig = plt.figure(figsize=(W / 72, H / 72), dpi=144)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W); ax.set_ylim(H, 0); ax.axis('off')
    ax.add_patch(plt.Rectangle((PL, PT_), PW, PH, facecolor='none',
                               edgecolor='#BFBFBF', lw=1.2, zorder=1))
    for i in range(n_div + 1):
        yy = PT_ + PH * i / n_div
        if 0 < i < n_div:
            ax.plot([PL, PL + PW], [yy] * 2, color='#EBEBEB', lw=1, zorder=1)
        lv = L_TOP - (L_TOP - L_BOT) * i / n_div
        rv = R_TOP - (R_TOP - R_BOT) * i / n_div
        ax.text(PL - 10, yy, f'{lv:,.0f}', ha='right', va='center', fontsize=19, color='#333')
        ax.text(PL + PW + 10, yy, f'{rv:g}', ha='left', va='center', fontsize=19, color='#333')
    # x축: 매년 09-30
    y0, y1 = int(ser[0][0][:4]), int(ser[-1][0][:4])
    for yr in range(y0, y1 + 1):
        d = f'{yr}-09-30'
        if _dnum(d) < T0 or _dnum(d) > T1:
            continue
        ax.plot([X(d)] * 2, [PT_ + PH, PT_ + PH + 7], color='#BFBFBF', lw=1.2)
        ax.text(X(d) + 6, PT_ + PH + 20, d, fontsize=19, color='#333',
                rotation=32, ha='right', va='top', rotation_mode='anchor')
    ax.plot([X(d) for d, _ in px], [YL(v) for _, v in px], color=C_PX, lw=2, zorder=3)
    ax.plot([X(d) for d, _, _ in ser], [YR(p) for _, p, _ in ser], color=C_PE, lw=2, zorder=3)
    fig.savefig(OUT / 's10_chart.png', facecolor='white')
    plt.close(fig)
    # 주석 박스 + 점선 인출선 = 네이티브 편집 개체 (2026-07-14 사용자 지시) —
    # 앵커(이미지 px 좌표)만 계산해 add() 로 전달 (원본 앵커: PER→2021-03-02, 가격→2019-09-02)
    d_b = next((d for d, _, _ in ser if d >= '2021-03-02'), ser[len(ser) // 2][0])
    p_b = dict((d, p) for d, p, _ in ser)[d_b]
    d_r = next((d for d, _ in px if d >= '2019-09-02'), px[len(px) // 3][0])
    v_r = dict(px)[d_r]
    annos = [
        # (라벨, 색hex, 박스중심 x/y, 인출선 시작 x/y, 인출선 끝 x/y)
        ('12개월 선행 PER\n( MSCI Korea, 우)', '2E5E9E',
         (960, 100), (960, 130), (X(d_b), YR(p_b) - 4)),
        ('Price Index\n( MSCI Korea, 좌)', 'A83232',
         (420, 265), (420, 300), (X(d_r), YL(v_r) - 4)),
    ]
    return OUT / 's10_chart.png', (W, H), annos


def add(prs, ctx, page_label='10'):
    end = ctx['asof']
    kr = ctx['valuation']['MXKR']
    pe_now = kr['series'][-1][1]
    y, m, d = end.split('-')
    sub = f'{y}.{int(m)}.{int(d)} 기준 MXKR PER(12FWD): {pe_now:.2f}'
    sl = slide_scaffold(prs, 'base_slide10.png', 'MSCI KR', end, page_label, subtitle=sub)
    add_text(sl, 100, 208, 1400, 34, '한국 가격지수 vs 밸류에이션', 24 * PT_PER_PX,
             '222222', bold=True, align=PP_ALIGN.CENTER)
    png, (w, h), annos = gen_chart(kr, end)
    IX, IY = 40, 252                              # 차트 이미지 원점 (슬라이드 px)
    sl.shapes.add_picture(str(png), EX(IX), E(IY), E(w), E(h))
    # 주석 = 네이티브: 점선 인출선(커넥터) 먼저(뒤), 테두리 텍스트박스 나중(앞)
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from .common import set_ko_font
    BXW, BXH = 210, 58                            # 주석 박스 크기 (px)
    for label, chex, (bcx, bcy), (lx0, ly0), (lx1, ly1) in annos:
        ln = sl.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, EX(IX + lx0), E(IY + ly0), EX(IX + lx1), E(IY + ly1))
        ln.line.color.rgb = RGBColor.from_string(chex)
        ln.line.width = Pt(2 * PT_PER_PX)
        pd = ln.line._get_or_add_ln().makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.line._get_or_add_ln().append(pd)
        bx = sl.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            EX(IX + bcx - BXW / 2), E(IY + bcy - BXH / 2), E(BXW), E(BXH))
        bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor.from_string('FFFFFF')
        bx.line.color.rgb = RGBColor.from_string(chex)
        bx.line.width = Pt(2 * PT_PER_PX)
        bx.shadow.inherit = False
        tf = bx.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, seg in enumerate(label.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = seg
            r.font.name = 'Pretendard'; r.font.size = Pt(round(19 * PT_PER_PX, 1))
            r.font.color.rgb = RGBColor.from_string(chex)
            set_ko_font(r.font, 'Pretendard')
    add_text(sl, 900, 1028, 660, 26, f'· 자료: {kdate(end)}, Bloomberg, 한국투자신탁운용',
             BODY_PT - 2, '787878', align=PP_ALIGN.RIGHT)
    return sl
