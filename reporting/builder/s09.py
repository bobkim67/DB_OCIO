"""슬라이드 9 — Total Return Breakdown: PER×EPS 분해 누적바 + 표 (이미지)."""
import math

from .common import (OUT, E, EX, add_text, slide_scaffold, plt, BODY_PT,
    PT_PER_PX, PP_ALIGN, kdate)

# (표시명, 심볼) — 발송본 14열. 데이터 없으면 빈칸.
CATS = [('전세계', 'MXWD'), ('미국', 'MXUS'), ('미국\n빅테크 7 Plus', 'BUBT7P'),
        ('미국\n성장주', 'MXUS000G'), ('미국\n가치주', 'MXUS000V'), ('DM ex US', 'MXWOU'),
        ('EM', 'MXEF'), ('한국', 'MXKR'), ('중국', 'MXCN'), ('일본', 'MXJP'),
        ('독일', 'MXDE'), ('영국', 'MXGB'), ('브라질', 'MXBR'), ('남아공', 'MXZA')]
RANK_SCOPE = [('DM ex US', 'MXWOU'), ('US Growth', 'MXUS000G'), ('US', 'MXUS'),
              ('Global', 'MXWD'), ('US Value', 'MXUS000V'), ('EM', 'MXEF'), ('Korea', 'MXKR')]
C_B, C_R, C_G = '#2C5A9C', '#A83232', '#6E8B3D'


def decomp(val, start, end):
    """구간 PER/EPS 변화율 분해: {sym: (per%, eps%, cross%, total%)}"""
    out = {}
    for sym, v in val.items():
        ser = [r for r in v['series'] if start <= r[0] <= end]
        if len(ser) < 2:
            continue
        # 구간 시작 = start 이후 첫 관측 (원본 동일)
        s0, s1 = ser[0], ser[-1]
        perc = (s1[1] / s0[1] - 1) * 100
        epsc = (s1[2] / s0[2] - 1) * 100
        tot = ((1 + perc / 100) * (1 + epsc / 100) - 1) * 100
        out[sym] = (round(perc, 1), round(epsc, 1), round(tot - perc - epsc, 1), round(tot, 1))
    return out


def subtitle_rank(dec):
    ranked = sorted((nm for nm, sym in RANK_SCOPE if sym in dec),
                    key=lambda nm: dec[dict(RANK_SCOPE)[nm]][3])
    return ' < '.join(ranked[:-1]) + ' <<< ' + ranked[-1] if len(ranked) > 1 else ''


def gen_chart(dec):
    """누적바 차트만 이미지로 (표는 네이티브 — 2026-07-14 사용자 지시)."""
    W, H = 1520, 428
    PL, PR, PT_, PH = 112, 17, 14, 400
    fig = plt.figure(figsize=(W / 72, H / 72), dpi=144)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W); ax.set_ylim(H, 0); ax.axis('off')
    vals = [v for sym in (s for _, s in CATS) if sym in dec for v in dec[sym][:3]]
    pos_max = max(sum(max(v, 0) for v in dec[s][:3]) for _, s in CATS if s in dec)
    neg_min = min(sum(min(v, 0) for v in dec[s][:3]) for _, s in CATS if s in dec)
    YMAX = max(100, math.ceil(pos_max / 50) * 50)
    YMIN = min(-50, math.floor(neg_min / 50) * 50)
    S = PH / (YMAX - YMIN)
    y0 = PT_ + YMAX * S
    PW = W - PL - PR
    for v in range(YMAX, YMIN - 1, -50):
        yy = PT_ + (YMAX - v) * S
        ax.plot([PL, PL + PW], [yy] * 2, color='#7f7f7f' if v == 0 else '#DDDDDD',
                lw=2 if v == 0 else 1, zorder=1)
        ax.text(PL - 10, yy, f'{v}%', ha='right', va='center', fontsize=16, color='#333')
    slot = PW / len(CATS)
    BW = 56
    for i, (nm, sym) in enumerate(CATS):
        cx = PL + slot * (i + 0.5)
        if sym not in dec:
            continue
        up = dn = 0.0
        for col, v in zip((C_B, C_R, C_G), dec[sym][:3]):
            if v is None:
                continue
            h = abs(v) * S
            if v >= 0:
                top = y0 - (up + v) * S; up += v
            else:
                top = y0 + dn * S; dn += -v
            ax.add_patch(plt.Rectangle((cx - BW / 2, top), BW, h, facecolor=col,
                                       edgecolor='none', zorder=3))
    fig.savefig(OUT / 's9_chart.png', facecolor='white')
    plt.close(fig)
    return OUT / 's9_chart.png', (W, H)


def add_native_table(sl, dec, px_x, px_y):
    """차트 하단 데이터 표 — 네이티브 편집 표 (2026-07-14 사용자 지시).

    구 이미지 표와 동일 지오메트리: 라벨열 112px + 14열, 헤더 52px + 4행 44px,
    전체 그리드 #C9C9C9, 총수익행 강조, PER/EPS/기타 라벨에 색상 스와치(■).
    """
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
    from .common import E, EX, EMU_PER_PX, PT_PER_PX, set_ko_font

    PW = 1520 - 112 - 17
    col_w = [112] + [PW / len(CATS)] * len(CATS)
    row_h = [52, 44, 44, 44, 44]
    labels = ['총수익', 'PER 변화율', 'EPS 변화율', '기타']
    swatch = {1: C_B, 2: C_R, 3: C_G}
    FS = Pt(round(15 * PT_PER_PX, 1))          # 이미지 표 fontsize 15px 등가
    FS_H = Pt(round(14.5 * PT_PER_PX, 1))

    gf = sl.shapes.add_table(5, len(CATS) + 1, EX(px_x), E(px_y),
                             E(sum(col_w)), E(sum(row_h)))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is not None:
        sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'   # No Style, No Grid
    for i, w in enumerate(col_w):
        tbl.columns[i].width = E(w)
    for i, h in enumerate(row_h):
        tbl.rows[i].height = E(h)

    def _borders(cell):
        # tcPr 자식은 스키마 시퀀스(lnL→lnR→lnT→lnB→…→fill) 순서 필수 —
        # 역순 삽입 시 PowerPoint 가 파일 자체를 거부 (2026-07-14 확인)
        tcPr = cell._tc.get_or_add_tcPr()
        for idx, tag in enumerate(('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB')):
            ln = tcPr.makeelement(qn(tag), {'w': '6350', 'cap': 'flat'})   # 0.5pt
            sf = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr')).set('val', 'C9C9C9')
            tcPr.insert(idx, ln)

    def _put(cell, runs, fs, align, fill=None):
        """runs: [(text, bold, color_hex)] — '\\n' 은 문단 분리."""
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(fill)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string('FFFFFF')
        _borders(cell)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = False
        tf.margin_top = tf.margin_bottom = 0
        tf.margin_left = tf.margin_right = Emu(round(3 * EMU_PER_PX))
        def _prep(p):
            """문단 서식 — 문단당 1회만 (run 마다 재호출 금지: ea/cs 중복 → 파일 거부)."""
            p.alignment = align
            p.font.name = 'Pretendard'; p.font.size = fs       # 빈 셀 행높이 방지
            set_ko_font(p.font, 'Pretendard')

        p = None
        for text, bold, color in runs:
            for j, seg in enumerate(text.split('\n')):
                if p is None:
                    p = tf.paragraphs[0]; _prep(p)
                elif j > 0:
                    p = tf.add_paragraph(); _prep(p)
                if not seg:
                    continue
                r = p.add_run(); r.text = seg
                r.font.name = 'Pretendard'; r.font.size = fs
                r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
                set_ko_font(r.font, 'Pretendard')

    # 헤더행: 좌측 빈 셀 + 카테고리명 (2줄 헤더는 문단 분리)
    _put(tbl.cell(0, 0), [('', False, '222222')], FS_H, PP_ALIGN.CENTER)
    for i, (nm, _s) in enumerate(CATS):
        _put(tbl.cell(0, i + 1), [(nm, True, '222222')], FS_H, PP_ALIGN.CENTER,
             fill='F5F5F5')
    # 데이터행
    for ri, lab in enumerate(labels):
        row = ri + 1
        fill0 = 'F2F2F2' if ri == 0 else 'FAFAFA'
        if ri in swatch:
            _put(tbl.cell(row, 0), [('■ ', False, swatch[ri].lstrip('#')),
                                    (lab, True, '222222')], FS, PP_ALIGN.LEFT, fill=fill0)
        else:
            _put(tbl.cell(row, 0), [(lab, True, '222222')], FS, PP_ALIGN.LEFT, fill=fill0)
        for i, (nm, sym) in enumerate(CATS):
            fill = 'F2F2F2' if ri == 0 else None
            if sym in dec:
                v = dec[sym][3] if ri == 0 else dec[sym][ri - 1]
                _put(tbl.cell(row, i + 1), [(f'{v:.1f}%', ri == 0, '222222')],
                     FS, PP_ALIGN.CENTER, fill=fill)
            else:
                _put(tbl.cell(row, i + 1), [('', False, '222222')], FS,
                     PP_ALIGN.CENTER, fill=fill)
    return gf


def add(prs, ctx, page_label='9'):
    # s9~16 은 커스텀 구간과 무관하게 YTD 고정 (2026-07-13 사용자 확정)
    start = ctx.get('ytd_start') or ctx['series_ytd'][0]['date']
    end = ctx['asof']
    dec = decomp(ctx['valuation'], start, end)
    sl = slide_scaffold(prs, 'base_slide09.png', 'Total Return Breakdown', end,
                        page_label, subtitle=subtitle_rank(dec))
    title = f"글로벌 주식 총 수익률 분석(USD, {start.replace('-', '')} ~ {end.replace('-', '')})"
    add_text(sl, 100, 208, 1400, 34, title, 24 * PT_PER_PX, '222222', bold=True,
             align=PP_ALIGN.CENTER)
    png, (w, h) = gen_chart(dec)
    sl.shapes.add_picture(str(png), EX(40), E(252), E(w), E(h))
    add_native_table(sl, dec, 40, 252 + 428)      # 구 이미지 표와 동일 위치(y680)
    add_text(sl, 100, 1018, 900, 26,
             '총수익은 12M 선행 PER×EPS 기반 가격수익률(USD, 배당 제외), 기타는 PER·EPS 변동의 교차항',
             13 * PT_PER_PX * 1.4, '777777')
    add_text(sl, 900, 1018, 660, 26, f'· 자료: {kdate(end)}, Bloomberg, 한국투자신탁운용',
             BODY_PT - 2, '787878',
             align=PP_ALIGN.RIGHT)
    return sl
