"""슬라이드 4 — 연초 이후: 자산군별 기간수익률 표(그라디언트 바, 이미지) + 코멘트(편집 텍스트).

표 수치 = SCIP 라이브. v1=표시통화 수익률, v2=원화환산 = (1+v1)x(1+USDKRW수익률)-1
(발송본 역산 검증: S&P500 9.31→17.56, Gold -7.82→-0.86). KR 행은 v2=v1.
해외 행 윈도우 = T-1 (발송본: 해외 2025.12.30~2026.06.29). 코멘트/헤드라인 = 수동 영역
(원본 = slide04_data.json 수동 작성) → 네이티브 편집 텍스트, 기본값은 참조 JSON.
"""
import datetime
import json

import numpy as np
import pymysql

from .common import (OUT, ROOT, E, add_text, slide_scaffold, plt, BODY_PT,
                     PT_PER_PX, PP_ALIGN, kdate)
from modules.data_loader import parse_data_blob

# (type, cat, cat_span, sub, sub_span, label, bench, did, sids, cur)
#   type: grp(회색 요약)/krow(피치 국내)/row. did=None → 빈칸(–).
S4_ROWS = [
    ('grp',  '주식', 11, None, 0, '글로벌',   'MSCI ACWI',                334, (9,),        'USD'),
    ('krow', None, 0, '국내', 2, '시장',      'KOSPI',                    253, (9, 15),     'KRW'),
    ('krow', None, 0, None, 0, '일반',        'KOSPI 200',                225, (9, 15),     'KRW'),
    ('row',  None, 0, '미국', 5, '시장',      'S&P500',                   271, (9, 6),      'USD'),
    ('row',  None, 0, None, 0, '성장',        'CRSP US Large Cap Growth', 431, (9,),        'USD'),
    ('row',  None, 0, None, 0, '가치',        'CRSP US Large Cap Value',  433, (9,),        'USD'),
    ('row',  None, 0, None, 0, '중소형',      'Russell 2000',             338, (9,),        'USD'),
    ('row',  None, 0, None, 0, '고배당',      'DJ Dividend 100',          275, (9,),        'USD'),
    ('row',  None, 0, '미국외 선진국', 1, '일반', 'MSCI World ex US',     339, (9,),        'USD'),
    ('row',  None, 0, '선진국', 1, '일반',    'MSCI World',               453, (48, 9, 6),  'USD'),
    ('row',  None, 0, '신흥시장', 1, '일반',  'MSCI EM',                  340, (9,),        'USD'),
    ('grp',  '채권', 8, None, 0, '글로벌 (UH)', 'Barclays Global Agg.',    58, (39, 9),     'USD'),
    ('krow', None, 0, '국내', 3, '국채 3년',  '매경채권지수',              422, (9,),        'KRW'),
    ('krow', None, 0, None, 0, '국고 10년',   'KRX 10년국채지수',          421, (9,),        'KRW'),
    ('krow', None, 0, None, 0, '종합채권',    'KBP 종합지수',              257, (9,),        'KRW'),
    ('row',  None, 0, '미국', 3, '종합채권',  'Barclays US Agg.',          278, (9, 39, 6), 'USD'),
    # 발송본 수치의 실제 소스 = BBG Corporate TR(450/401 ds9) — 라벨은 발송본(IBOX) 유지
    ('row',  None, 0, None, 0, '투자등급',    'IBOX Investment Grade',     450, (9,),       'USD'),
    ('row',  None, 0, None, 0, '하이일드',    'IBOX High Yield',           401, (9,),       'USD'),
    ('row',  None, 0, '신흥시장', 1, '달러국채', 'JP Morgan EM Bond',      244, (9, 6, 15), 'USD'),
    ('row',  '대체투자', 5, '원자재', 2, 'WTI', 'WTI Crude Oil',            98, (15, 6),    'USD'),
    ('row',  None, 0, None, 0, 'Gold',        'Gold',                      277, (15, 6),    'USD'),
    ('row',  None, 0, '인프라', 1, '글로벌인프라', 'S&P Global Infrastructure', 40, (6, 15, 9), 'USD'),
    ('row',  None, 0, '부동산', 2, '미국리츠', 'MSCI US REITs',            317, (6, 15, 9), 'USD'),
    ('row',  None, 0, None, 0, '미국 제외 글로벌 리츠', 'S&P Global ex US Property', 39, (6, 15, 48, 9), 'USD'),
    ('grp',  '통화', 2, None, 0, '달러 인덱스', '달러 인덱스',              105, (48, 6),   'USD'),
    ('krow', None, 0, None, 0, 'USD/KRW',     'USDKRW',                    31, (6,),        'FX'),
]

G_POS = ('#6CC68B', '#E9F6EE', '#4EA36A')   # 그라디언트 시작/끝/테두리 (양수=녹색)
G_NEG = ('#FF5860', '#FFE9EA', '#D64550')


def _prev_bday(cur, did, sids, target):
    return None


def _load_series(cur, did, sids, end_date):
    """우선순위 dataseries 로 {date: value} — 첫 유효 계열."""
    for sid in sids:
        cur.execute(
            "SELECT DATE(timestamp_observation) d, data FROM back_datapoint "
            "WHERE dataset_id=%s AND dataseries_id=%s AND DATE(timestamp_observation)<=%s "
            "ORDER BY timestamp_observation", (did, sid, end_date))
        ser = {}
        for r in cur.fetchall():
            try:
                v = parse_data_blob(r['data'])
                if isinstance(v, dict):
                    v = v.get('USD') or v.get('KRW') or next(iter(v.values()))
                ser[r['d'].strftime('%Y-%m-%d')] = float(v)
            except Exception:
                continue
        if len(ser) > 50:
            return ser
    return {}


def compute_rows(end_date: str, start_date: str | None = None) -> dict:
    """{'rows': [발송본 slide04 rows 스키마], ...} — SCIP 라이브.

    start_date 지정 시 그 앵커(≤)부터, 미지정 시 전년말(YTD — 발송본 컨셉).
    """
    end = datetime.date.fromisoformat(end_date)
    kr_start_target = (datetime.date.fromisoformat(start_date) if start_date
                       else datetime.date(end.year - 1, 12, 31))
    conn = pymysql.connect(host='192.168.195.55', user='solution', password='Solution123!',
                           db='SCIP', charset='utf8mb4', cursorclass=pymysql.cursors.DictCursor)
    try:
        cur = conn.cursor()
        fx = _load_series(cur, 31, (6,), end_date)

        def window_ret(ser, start_t, end_t):
            ds = sorted(ser)
            s = next((d for d in reversed(ds) if d <= start_t.isoformat()), None)
            e = next((d for d in reversed(ds) if d <= end_t.isoformat()), None)
            if not s or not e or s == e:
                return None, None, None
            return ser[e] / ser[s] - 1, s, e

        fx_ret, kr_s, kr_e = window_ret(fx, kr_start_target, end)
        # 해외 윈도우 = 대시보드 규약(2026-07-13 사용자 확정): 시작·종료 모두
        #   국내 앵커의 전영업일(T-1). 예: 국내 12/31~6/30 → 해외 12/30~6/29.
        #   (발송본 계산은 시작 12/31 이었으나 라벨은 12/30 — 라벨 쪽이 규약)

        rows, bar_max = [], 0.0
        for (typ, cat, cspan, sub, sspan, label, bench, did, sids, curc) in S4_ROWS:
            v1 = v2 = None
            if did:
                ser = _load_series(cur, did, sids, end_date)
                if ser:
                    if curc == 'KRW':
                        r, _, _ = window_ret(ser, datetime.date.fromisoformat(kr_s), end)
                        v1 = v2 = r
                    elif curc == 'FX':
                        r, _, _ = window_ret(ser, datetime.date.fromisoformat(kr_s), end)
                        v1, v2 = r, None
                    else:                       # USD 표시 — 시작·종료 = 국내 앵커 T-1 + 원화환산
                        ds_ser = sorted(ser)
                        s_t = next((d for d in reversed(ds_ser) if d < kr_s), None)
                        e_t = next((d for d in reversed(ds_ser) if d < kr_e), None)
                        r = (ser[e_t] / ser[s_t] - 1) if s_t and e_t and s_t != e_t else None
                        v1 = r
                        v2 = ((1 + r) * (1 + fx_ret) - 1) if r is not None else None
            row = {'type': typ, 'label': label, 'bench': bench,
                   'v1': None if v1 is None else round(v1 * 100, 2),
                   'v2': None if v2 is None else round(v2 * 100, 2)}
            if cat:
                row['cat'] = cat; row['cat_span'] = cspan
            if sub:
                row['sub'] = sub; row['sub_span'] = sspan
            rows.append(row)
            for v in (row['v1'], row['v2']):
                if v is not None:
                    bar_max = max(bar_max, abs(v))
        fxd = sorted(fx)
        os_s = next((d for d in reversed(fxd) if d < kr_s), kr_s)
        os_e = next((d for d in reversed(fxd) if d < kr_e), kr_e)
    finally:
        conn.close()
    period_line = (f"한국: {kr_s.replace('-', '.')} ~ {kr_e.replace('-', '.')}, "
                   f"해외: {os_s.replace('-', '.')} ~ {os_e.replace('-', '.')}")
    return {'rows': rows, 'bar_max': bar_max, 'period_line': period_line,
            'kr_start': kr_s}


def _grad(ax, x, y, w, h, neg):
    c0, c1, ec = G_NEG if neg else G_POS
    from matplotlib.colors import LinearSegmentedColormap
    cm = LinearSegmentedColormap.from_list('g', [c0, c1])
    ax.imshow(np.linspace(0, 1, 64).reshape(1, -1), cmap=cm, aspect='auto',
              extent=[x, x + w, y + h, y], zorder=2)
    ax.add_patch(plt.Rectangle((x, y), w, h, facecolor='none', edgecolor=ec,
                               lw=0.8, zorder=3))


def gen_table(data):
    rows, bar_max = data['rows'], data['bar_max'] or 1.0
    COLS = [80, 100, 150, 105, 105, 320]
    W = sum(COLS) + 4
    TH, TD = 42, 27
    H = TH + len(rows) * TD + 4
    fig = plt.figure(figsize=(W / 72, H / 72), dpi=144)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W); ax.set_ylim(H, 0); ax.axis('off')
    FS = 15
    xs = [2]
    for c in COLS:
        xs.append(xs[-1] + c)

    def cell(x0, x1, y, h, text='', fill=None, bold=False, fs=FS, color='#000'):
        ax.add_patch(plt.Rectangle((x0, y), x1 - x0, h, facecolor=fill or 'white',
                                   edgecolor='#B7B7B7', lw=0.8, zorder=1))
        if text:
            ax.text((x0 + x1) / 2, y + h / 2, text, ha='center', va='center',
                    fontsize=fs, color=color, fontweight='bold' if bold else 'normal',
                    zorder=4)

    def bar_cell(x0, x1, y, v, krow, override=None):
        cell(x0, x1, y, TD, '', fill='#F8E6D6' if krow else 'white')
        if v is None:
            ax.text((x0 + x1) / 2, y + TD / 2, '–', ha='center', va='center',
                    fontsize=FS, color='#7f7f7f', zorder=4)
            return
        blen = abs(v) / bar_max * (x1 - x0 - 4)
        _grad(ax, x0 + 1, y + 3, max(blen, 1), TD - 6, v < 0)
        ax.text((x0 + x1) / 2, y + TD / 2, override or f'{v:.2f}%', ha='center',
                va='center', fontsize=FS, zorder=4)

    # 헤더
    cell(xs[0], xs[3], 2, TH, '자산군', fill='#5B9BD5', bold=True, color='white')
    cell(xs[3], xs[4], 2, TH, '기간 수익률\n(표시 통화)', fill='#5B9BD5', bold=True,
         color='white', fs=13)
    cell(xs[4], xs[5], 2, TH, '기간 수익률\n(원화환산)', fill='#5B9BD5', bold=True,
         color='white', fs=13)
    cell(xs[5], xs[6], 2, TH, '벤치마크', fill='#5B9BD5', bold=True, color='white')

    y = 2 + TH
    sub_cover = 0     # 직전 sub 셀 rowspan 이 이 행을 덮는지 (잔여 행수)
    for r in rows:
        krow = r['type'] == 'krow'
        grp = r['type'] == 'grp'
        # cat 셀 (rowspan)
        if 'cat' in r:
            cell(xs[0], xs[1], y, TD * r['cat_span'], r['cat'], bold=True)
        # sub / label
        if grp:
            cell(xs[1], xs[3], y, TD, r['label'], fill='#F2F2F2', bold=True)
            sub_cover = 0
        else:
            if 'sub' in r:
                cell(xs[1], xs[2], y, TD * r['sub_span'], r['sub'],
                     fill='#F8E6D6' if krow else 'white')
                sub_cover = r['sub_span']
            if sub_cover <= 0:
                # sub 셀이 없는 단독행(USD/KRW): 좌측 칸과 병합 — 윗행 달러인덱스와 동일
                # 전폭 셀이라 하단 테두리도 생김 (2026-07-14 사용자 지시)
                cell(xs[1], xs[3], y, TD, r['label'], fill='#F8E6D6' if krow else 'white')
            else:
                cell(xs[2], xs[3], y, TD, r['label'], fill='#F8E6D6' if krow else 'white')
            sub_cover -= 1
        v2_override = '–' if (r['v2'] is None and r['v1'] is not None) else None
        bar_cell(xs[3], xs[4], y, r['v1'], krow)
        bar_cell(xs[4], xs[5], y, r['v2'], krow, v2_override if r['v2'] is None else None)
        cell(xs[5], xs[6], y, TD, r['bench'], fill='#F2F2F2' if grp else
             ('#F8E6D6' if krow else 'white'), bold=grp, fs=13.5)
        y += TD
    p = OUT / 's4_table.png'
    fig.savefig(p, facecolor='white')
    plt.close(fig)
    return p, (W, H)


def _default_manual():
    """코멘트/헤드라인 기본값 — 발송본 slide04_data.json (수동 작성 정본)."""
    try:
        j = json.loads((ROOT / 'reference' / 'slide04_data.json').read_text(encoding='utf-8'))
        return {'headline': j.get('headline', ''), 'comments': j.get('comments', [])}
    except Exception:
        return {'headline': '', 'comments': []}


def add(prs, ctx, page_label='4'):
    end = ctx['asof']
    data = compute_rows(end, None if ctx.get('is_ytd', True) else ctx.get('period_start'))
    # 코멘트 자동 생성 (순위=표 유도 + 정성=승인 시장코멘트 LLM 요약, 캐시/수동편집 가능)
    manual = ctx.get('s4_manual')
    if manual is None:
        from .s04_comment import build_manual
        try:
            manual = build_manual(data, end,
                                  tag='' if ctx.get('is_ytd', True) else ctx.get('period_start', ''))
        except Exception as e:            # noqa: BLE001
            print(f'[s4] 코멘트 자동생성 실패 → 참조 JSON 폴백: {e}')
            manual = _default_manual()
    y_, m_, d_ = data['kr_start'].split('-')
    title = (f'연초({y_}.{m_}.{d_}) 이후' if ctx.get('is_ytd', True)
             else f'{y_}.{m_}.{d_} 이후')
    sl = slide_scaffold(prs, 'base_slide07.png', title, end,
                        page_label, subtitle=data['period_line'])
    # 헤드라인 (편집 텍스트)
    if manual['headline']:
        add_text(sl, 52, 210, 1500, 30, manual['headline'], BODY_PT + 1, '111111', bold=True)
    # 좌: 자산군표 (그라디언트 바 — 이미지)
    png, (w, h) = gen_table(data)
    sl.shapes.add_picture(str(png), E(40), E(255), E(w), E(h))
    # 우: 코멘트 = 네이티브 표 (2026-07-13 사용자 지정 — 2열, 가로 구분선만,
    #   라벨=가운데/중간, 상세=왼쪽/중간, 글머리(·) hanging indent, 총높이=좌측 표와 동일)
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
    from .common import EMU_PER_PX, set_ko_font

    cx, cw = 40 + w + 30, 1600 - (40 + w + 30) - 30
    secs = manual['comments']
    # 행높이: 콘텐츠 줄수 비례 배분, 합 = 좌측 표 이미지 높이 h
    est = [sum(max(1, -(-len(s) // 30)) for s in sec['lines']) + 0.6 for sec in secs]
    tot = sum(est)
    row_h = [max(60, round(h * e / tot)) for e in est]
    row_h[-1] = max(60, h - sum(row_h[:-1]))

    gf = sl.shapes.add_table(len(secs), 2, E(cx), E(255), E(cw), E(h))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is not None:
        sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'   # No Style, No Grid
    tbl.columns[0].width = E(120)
    tbl.columns[1].width = E(cw - 120)
    for i, rh in enumerate(row_h):
        tbl.rows[i].height = E(rh)

    def _hline(cell, tag):                 # 섹션 가로 구분선 (E4E4E4)
        tcPr = cell._tc.get_or_add_tcPr()
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set('w', '9525'); ln.set('cap', 'flat')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr')).set('val', 'E4E4E4')
        tcPr.insert(0, ln)

    for ri, sec in enumerate(secs):
        for ci in range(2):
            cell = tbl.cell(ri, ci)
            cell.fill.background()                       # 투명 (줄무늬 없음)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE     # 세로 중간
            if ri == 0:
                _hline(cell, 'a:lnT')
            _hline(cell, 'a:lnB')
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = Emu(round(6 * EMU_PER_PX))
            tf.margin_left = Emu(round((4 if ci == 0 else 10) * EMU_PER_PX))
            tf.margin_right = Emu(round(8 * EMU_PER_PX))
        # 라벨 (가운데 정렬)
        lc = tbl.cell(ri, 0).text_frame
        for i, seg in enumerate(sec['label'].split('<br>')):
            p = lc.paragraphs[0] if i == 0 else lc.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = seg
            r.font.name = 'Pretendard'; r.font.size = Pt(BODY_PT + 1)
            r.font.bold = True; r.font.color.rgb = RGBColor.from_string('111111')
            set_ko_font(r.font, 'Pretendard')
        # 상세 (왼쪽 정렬 + · 글머리 hanging indent — 줄바꿈 시 자동 들여쓰기)
        dc = tbl.cell(ri, 1).text_frame
        for i, line in enumerate(sec['lines']):
            p = dc.paragraphs[0] if i == 0 else dc.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            if i > 0:
                p.space_before = Pt(4)
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', '152400'); pPr.set('indent', '-152400')
            pPr.append(pPr.makeelement(qn('a:buFont'), {'typeface': 'Pretendard'}))
            pPr.append(pPr.makeelement(qn('a:buChar'), {'char': '·'}))
            r = p.add_run(); r.text = line.lstrip('· ').strip()
            r.font.name = 'Pretendard'; r.font.size = Pt(BODY_PT)
            r.font.color.rgb = RGBColor.from_string('1A1A1A')
            set_ko_font(r.font, 'Pretendard')
    return sl


def check(end_date: str):
    """발송본 slide04_data.json 과 라이브 수치 대조 (매핑 검증)."""
    data = compute_rows(end_date)
    ref = json.loads((ROOT / 'reference' / 'slide04_data.json').read_text(encoding='utf-8'))
    print(f"period: live '{data['period_line']}' vs ref '{ref['period_line']}'")
    print(f"{'label':26s} {'ref v1':>8s} {'live v1':>8s} {'ref v2':>8s} {'live v2':>8s}")
    for rr, lr in zip(ref['rows'], data['rows']):
        f = lambda v: '   –' if v is None else f'{v:7.2f}'
        mark = ''
        if rr['v1'] is not None and lr['v1'] is not None and abs(rr['v1'] - lr['v1']) > 0.15:
            mark = '  <<<'
        print(f"{(rr.get('sub') or '') + ' ' + rr['label']:26s} {f(rr['v1'])} {f(lr['v1'])} "
              f"{f(rr['v2'])} {f(lr['v2'])}{mark}")
