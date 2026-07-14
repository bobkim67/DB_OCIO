"""슬라이드 16 — 경기국면: 불릿 + 국면정의표 + 스파이럴(SCIP 원천데이터 matplotlib 직접 렌더)."""
import numpy as np
from matplotlib.collections import LineCollection
from matplotlib.colors import LinearSegmentedColormap

from .common import (
    OUT, BODY_PT, HDR_BLUE, Z1, Z2, INK, RED, BLUE,
    sv, E, EX, add_text, add_bullets, add_table, slide_scaffold, plt,
)

# 스파이럴 팔레트 (SCIP 다크 테마 근사)
BG = '#10141F'
QUAD_TINT = {'TL': '#131A2A', 'TR': '#12201C', 'BL': '#1E1420', 'BR': '#201C12'}
QUAD_LABEL = [('회복', '#4FA3F7', 'TL'), ('팽창', '#35D0A0', 'TR'),
              ('침체', '#FF6B6B', 'BL'), ('둔화', '#E8C24A', 'BR')]


def gen_spiral(ctx):
    """G7 X_spiral/Y_spiral 경로 — 데이터 범위 줌인 + 월별 점 + 시간 그라디언트.

    2026-07-13 개선: ① 대칭 전체뷰 → 데이터 bbox 줌인 ② 매월 점(반기는 크게+라벨)
    ③ 현재 배지는 이미지에서 제거 — pptx 네이티브 도형/텍스트로 (badge_px 반환).
    """
    rg = ctx['regime']
    rows = rg['rows']
    x = np.array([r['X_spiral'] for r in rows])
    y = np.array([r['Y_spiral'] for r in rows])
    W, H = 930, 790
    # 데이터 bbox + 12% 마진 (0축은 항상 포함해 4분면 유지)
    xmin, xmax = min(x.min(), 0), max(x.max(), 0)
    ymin, ymax = min(y.min(), 0), max(y.max(), 0)
    mx, my = (xmax - xmin) * 0.12, (ymax - ymin) * 0.14
    xmin, xmax, ymin, ymax = xmin - mx, xmax + mx, ymin - my, ymax + my
    fig = plt.figure(figsize=(W / 72, H / 72), dpi=144)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(xmin, xmax); ax.set_ylim(ymin, ymax)
    ax.set_facecolor(BG); fig.patch.set_facecolor(BG)
    ax.axis('off')
    # 4분면 틴트 + 십자축 (뷰 범위 기준)
    ax.add_patch(plt.Rectangle((xmin, 0), -xmin, ymax, color=QUAD_TINT['TL'], zorder=0))
    ax.add_patch(plt.Rectangle((0, 0), xmax, ymax, color=QUAD_TINT['TR'], zorder=0))
    ax.add_patch(plt.Rectangle((xmin, ymin), -xmin, -ymin, color=QUAD_TINT['BL'], zorder=0))
    ax.add_patch(plt.Rectangle((0, ymin), xmax, -ymin, color=QUAD_TINT['BR'], zorder=0))
    ax.axhline(0, color='#2A3040', lw=1.2, zorder=1)
    ax.axvline(0, color='#2A3040', lw=1.2, zorder=1)
    for name, col, q in QUAD_LABEL:
        px = xmin + (xmax - xmin) * (0.03 if q[1] == 'L' else 0.97)
        py = ymin + (ymax - ymin) * (0.96 if q[0] == 'T' else 0.04)
        ax.text(px, py, name, color=col, fontsize=30, fontweight='bold',
                ha='left' if q[1] == 'L' else 'right',
                va='top' if q[0] == 'T' else 'bottom', zorder=5)
    # 경로: 시간 그라디언트 (cyan → magenta)
    cmap = LinearSegmentedColormap.from_list('spiral', ['#22C4E6', '#5A8CF0', '#B05CF0', '#E93FD4'])
    pts = np.array([x, y]).T.reshape(-1, 1, 2)
    segs = np.concatenate([pts[:-1], pts[1:]], axis=1)
    lc = LineCollection(segs, cmap=cmap, array=np.linspace(0, 1, len(segs)),
                        linewidths=3.4, zorder=3)
    ax.add_collection(lc)
    # 월별 점 (매월 작게), 반기(05/11)는 크게 + 라벨
    for i, r in enumerate(rows):
        mm = r['TIME_PERIOD'][5:7]
        if mm in ('05', '11'):
            ax.plot(x[i], y[i], 'o', ms=6.5, color='#E8473B', zorder=4)
            ax.annotate(r['TIME_PERIOD'], (x[i], y[i]), textcoords='offset points',
                        xytext=(9, 7), fontsize=13.5, color='#9AA2B4', zorder=4)
        else:
            ax.plot(x[i], y[i], 'o', ms=3.6, color='#E8473B', alpha=0.75, zorder=4)
    # 현재 지점 (노란 링만 — 배지는 네이티브)
    ax.plot(x[-1], y[-1], 'o', ms=13, mfc='#F5D46A', mec='#FFF0B0', mew=2, zorder=6)
    p = OUT / 's16_spiral.png'
    fig.savefig(p, facecolor=BG)
    plt.close(fig)
    badge_px = ((x[-1] - xmin) / (xmax - xmin) * W,
                (ymax - y[-1]) / (ymax - ymin) * H,
                bool(x[-1] > (xmin + xmax) / 2))
    return p, (W, H), badge_px


def add(prs, ctx, page_label='16'):
    rg = ctx['regime']
    sl = slide_scaffold(prs, 'base_slide07.png', '경기국면', ctx['asof'], page_label)
    W = 'FFFFFF'; K = INK

    # 좌: 불릿 (기본 문구 자동 생성 — PPT 에서 편집)
    m = int(rg['report_ym'][5:7])
    bullets = [
        (f'· 한국투자신탁운용 솔루션전략부의 경기국면모형에 따르면, {m}월 현재 '
         f'{rg["latest_phase"]}국면에 머물고 있는 것으로 판단'),
        '· 비침체 포트폴리오 유지',
        '· 국내주식은 12선행 EPS기준 7~8배 수준으로 비중 확대 유지',
    ]
    add_bullets(sl, 52, 240, 520, bullets, pt_size=BODY_PT + 1)

    # 좌하: 국면 정의표 (감소=빨강 / 증가=파랑)
    rows = [
        (HDR_BLUE, [('국면', True, W), ('밸류에이션', True, W), ('기업 이익', True, W)]),
        (Z1, [('둔화', False, K), ('감소', True, RED), ('증가', True, BLUE)]),
        (Z2, [('침체', False, K), ('감소', True, RED), ('감소', True, RED)]),
        (Z1, [('회복', False, K), ('증가', True, BLUE), ('감소', True, RED)]),
        (Z2, [('팽창', False, K), ('증가', True, BLUE), ('증가', True, BLUE)]),
    ]
    add_table(sl, 52, 620, [160, 200, 200], [sv(40)] + [sv(38)] * 4, rows)

    # 우: 스파이럴 (SCIP 원천데이터 렌더) + 현재 배지 = 네이티브 편집 도형 (2026-07-13)
    png, (w, h), (bx, by, right_side) = gen_spiral(ctx)
    x0, y0 = 620, 230
    sl.shapes.add_picture(str(png), EX(x0), E(y0), E(w), E(h))
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN as _AL, MSO_ANCHOR as _AN
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from .common import set_ko_font
    BW, BH = 190, 66
    bx_s = x0 + bx + (-BW - 26 if right_side else 26)          # 우측 끝이면 왼쪽에 배치
    by_s = y0 + by + 20
    bx_s = max(x0 + 8, min(bx_s, x0 + w - BW - 8))
    by_s = max(y0 + 8, min(by_s, y0 + h - BH - 8))
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, EX(bx_s), E(by_s), E(BW), E(BH))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string('10141F')
    sh.line.color.rgb = RGBColor.from_string('F5D46A'); sh.line.width = Pt(1.2)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = _AN.MIDDLE
    for i, line in enumerate([f"{rg['report_ym'].replace('-', '.')}  현재",
                              f"{rg['latest_phase']}국면"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = _AL.CENTER
        r = p.add_run(); r.text = line
        r.font.name = 'Pretendard'; r.font.size = Pt(11)
        r.font.bold = True; r.font.color.rgb = RGBColor.from_string('F5D46A')
        set_ko_font(r.font, 'Pretendard')
    return sl
