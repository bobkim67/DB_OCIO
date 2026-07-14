"""정적 슬라이드 — 표지·목차·섹션표지 (사용자 재구성 202606 편집본에서 이식, 2026-07-14).

스펙 = reference/static_slides_202606.json (열린 편집본 COM 추출 — pt 좌표·런 단위 서식).
그림 = template/static/*.png (클립보드 4배 캡처). 그룹은 절대좌표로 평면화해 그림
(사용자가 PPT 에서 재그룹 가능). 표지 년월은 ctx['asof'] 로 동적 치환.
"""
import json

from .common import ROOT, PTE, set_ko_font

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

SPEC_PATH = ROOT / 'reference' / 'static_slides_202606.json'
STATIC_DIR = ROOT / 'template' / 'static'

# 그림 자산 매핑 (COM 도형명 → 추출 PNG)
_PIC_ASSETS = {
    'Picture 8': 'cover_logo_top.png',       # 표지 상단 로고
    'Picture 26': 'cover_icon_pension.png',  # 표지 연금 컨텐츠 아이콘 (Group 25 내)
    'Picture 1': 'footer_logo_band.png',     # 하단 로고띠 (목차)
    'Picture 7': 'footer_logo_band.png',     # 하단 로고띠 (섹션 — 슬라이드별 이름 상이)
    'Picture 10': 'footer_logo_band.png',
    'Picture 9': 'footer_logo_band.png',
    'Picture 3': 'footer_logo_band.png',
}

_ALIGN = {1: PP_ALIGN.LEFT, 2: PP_ALIGN.CENTER, 3: PP_ALIGN.RIGHT, 4: PP_ALIGN.JUSTIFY}
_ANCHOR = {1: MSO_ANCHOR.TOP, 3: MSO_ANCHOR.MIDDLE, 4: MSO_ANCHOR.BOTTOM}

_spec_cache = None


def _spec():
    global _spec_cache
    if _spec_cache is None:
        _spec_cache = json.loads(SPEC_PATH.read_text(encoding='utf-8'))
    return _spec_cache


def _fill_text(sh, sp, subs):
    """스펙의 문단/런을 도형 텍스트프레임에 채움."""
    tf = sh.text_frame
    tf.word_wrap = sp.get('wordWrap', 0) == -1
    tf.vertical_anchor = _ANCHOR.get(sp.get('vAnchor', 1), MSO_ANCHOR.TOP)
    m = sp.get('margins')
    if m:
        tf.margin_left, tf.margin_top, tf.margin_right, tf.margin_bottom = (
            PTE(m[0]), PTE(m[1]), PTE(m[2]), PTE(m[3]))
    for i, para in enumerate(sp.get('paras', [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = _ALIGN.get(para.get('align', 1), PP_ALIGN.LEFT)
        if para.get('spaceBefore'):
            p.space_before = Pt(para['spaceBefore'])
        if para.get('spaceWithin') and para['spaceWithin'] != 1:
            p.line_spacing = para['spaceWithin']
        for run in para.get('runs', []):
            text = run.get('text', '')
            for old, new in (subs or {}).items():
                text = text.replace(old, new)
            # 수직탭(\x0b) = PPT 소프트 줄바꿈 — 문단 분리로 재현 (a:t 에 그대로 넣으면
            # _x000B_ 문자로 노출됨, 2026-07-14 표지 QR 캡션에서 확인)
            for k, seg in enumerate(text.split('\x0b')):
                if k > 0:
                    p = tf.add_paragraph()
                    p.alignment = _ALIGN.get(para.get('align', 1), PP_ALIGN.LEFT)
                    if para.get('spaceWithin') and para['spaceWithin'] != 1:
                        p.line_spacing = para['spaceWithin']
                r = p.add_run(); r.text = seg
                r.font.name = run.get('font', 'Pretendard')
                r.font.size = Pt(run.get('size', 10))
                r.font.bold = run.get('bold', 0) == -1
                r.font.italic = run.get('italic', 0) == -1
                r.font.color.rgb = RGBColor.from_string(run.get('color', '000000'))
                set_ko_font(r.font, r.font.name)


def _apply_fill_line(sh, sp):
    fill = sp.get('fill')
    if fill and fill != 'none':
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string(fill)
        tr = sp.get('fillTransparency') or 0
        if tr > 0:                                  # 투명도 → a:alpha
            srgb = sh.fill.fore_color._xFill.find(qn('a:srgbClr'))
            if srgb is not None:
                a = srgb.makeelement(qn('a:alpha'), {'val': str(round((1 - tr) * 100000))})
                srgb.append(a)
    else:
        sh.fill.background()
    line = sp.get('line')
    if line and line != 'none':
        sh.line.color.rgb = RGBColor.from_string(line)
        sh.line.width = Pt(sp.get('lineW', 0.75))
        if sp.get('lineDash', 1) != 1:
            ln = sh.line._get_or_add_ln()
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False


def _draw_shape(sl, sp, subs):
    typ = sp['type']
    L, T, W, H = sp['L'], sp['T'], sp['W'], sp['H']
    if typ == 6:                                    # 그룹 → 절대좌표 평면화
        for c in sp.get('children', []):
            _draw_shape(sl, c, subs)
        return
    if typ == 13:                                   # 그림
        asset = _PIC_ASSETS.get(sp['name'])
        if asset and (STATIC_DIR / asset).exists():
            sl.shapes.add_picture(str(STATIC_DIR / asset), PTE(L), PTE(T), PTE(W), PTE(H))
        else:
            print(f"[s_static] 그림 자산 없음: {sp['name']} → 생략")
        return
    if typ == 9 or (typ == 1 and sp.get('autoShapeType') == -2):   # 선/커넥터
        ln = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     PTE(L), PTE(T), PTE(L + W), PTE(T + H))
        line = sp.get('line')
        if line and line != 'none':
            ln.line.color.rgb = RGBColor.from_string(line)
            ln.line.width = Pt(sp.get('lineW', 0.75))
        if sp.get('lineBeginArrow', 1) > 1 or sp.get('lineEndArrow', 1) > 1:
            lnEl = ln.line._get_or_add_ln()
            if sp.get('lineBeginArrow', 1) > 1:
                lnEl.append(lnEl.makeelement(qn('a:headEnd'), {'type': 'arrow'}))
            if sp.get('lineEndArrow', 1) > 1:
                lnEl.append(lnEl.makeelement(qn('a:tailEnd'), {'type': 'arrow'}))
        ln.shadow.inherit = False
        return
    if typ == 17:                                   # 텍스트박스
        sh = sl.shapes.add_textbox(PTE(L), PTE(T), PTE(W), PTE(H))
        _fill_text(sh, sp, subs)
        return
    if typ == 1:                                    # 오토셰이프 (rect=1 / rounded=5)
        mso = MSO_SHAPE.ROUNDED_RECTANGLE if sp.get('autoShapeType') == 5 else MSO_SHAPE.RECTANGLE
        sh = sl.shapes.add_shape(mso, PTE(L), PTE(T), PTE(W), PTE(H))
        adj = sp.get('adjustments')
        if adj and sp.get('autoShapeType') == 5:
            try:
                sh.adjustments[0] = adj[0]
            except Exception:                       # noqa: BLE001 — 조정값 실패는 무시
                pass
        _apply_fill_line(sh, sp)
        if sp.get('paras'):
            _fill_text(sh, sp, subs)
        if sp.get('rot'):
            sh.rotation = sp['rot']
        return
    print(f"[s_static] 미지원 도형 type={typ} '{sp['name']}' → 생략")


def _add_from_spec(prs, key, layout_name, ctx=None):
    from .common import layout_by_name
    spec = _spec()[key]
    sl = prs.slides.add_slide(layout_by_name(prs, layout_name))
    subs = {}
    if ctx:
        y, m = ctx['asof'][:4], int(ctx['asof'][5:7])
        subs['2026년 6월'] = f'{y}년 {m}월'          # 표지 년월 동적 치환
    for sp in spec['shapes']:
        _draw_shape(sl, sp, subs)
    return sl


def add_cover(prs, ctx):
    return _add_from_spec(prs, 'slide1', '제목 및 내용', ctx)


def add_toc(prs, ctx=None):
    return _add_from_spec(prs, 'slide2', '제목 및 내용')


def add_section(prs, which):
    """which: 'slide3'(01 금융시장 리뷰) | 'slide5'(02 운용 및 성과 리뷰) |
    'slide8'(03 펀더멘털 점검) | 'slide15'(04 운용 계획)"""
    return _add_from_spec(prs, which, '1_Title and Content')
