"""슬라이드 13 — 미국 성장주와 금의 분산투자효과: 누적로그수익률 금 vs 미국성장주 (KRW).

데이터: SCIP — 금 XAU(277/ds15 FG Price), 미국성장주 MXUS000G(336/ds9 BBG TR),
USDKRW(31/ds6 blob 'USD'). 전부 1999-12-31~ 커버 (2026-07-13 검증).
"""
import datetime
import math

import pymysql

from .common import (OUT, E, EX, add_text, slide_scaffold, plt, BODY_PT,
                     PT_PER_PX, PP_ALIGN, kdate)
from modules.data_loader import parse_data_blob

C_GOLD, C_GROWTH = '#5B9BD5', '#ED7D31'
EVENTS = [('2000-10-02', 'IT 버블'), ('2008-10-01', '글로벌 금융위기'),
          ('2011-09-01', '유럽 재정위기'), ('2020-03-16', '팬데믹'),
          ('2022-06-15', '미국 금리인상'), ('2026-06-16', '미-이란전쟁')]


def _dnum(d):
    y, m, dd = map(int, d.split('-'))
    return datetime.date(y, m, dd).toordinal()


def _load(end_date):
    conn = pymysql.connect(host='192.168.195.55', user='solution', password='Solution123!',
                           db='SCIP', charset='utf8mb4', cursorclass=pymysql.cursors.DictCursor)
    try:
        cur = conn.cursor()
        out = {}
        for key, did, sid in [('gold', 277, 15), ('growth', 336, 9), ('fx', 31, 6)]:
            cur.execute(
                "SELECT DATE(timestamp_observation) d, data FROM back_datapoint "
                "WHERE dataset_id=%s AND dataseries_id=%s AND DATE(timestamp_observation)<=%s "
                "ORDER BY timestamp_observation", (did, sid, end_date))
            ser = {}
            for r in cur.fetchall():
                try:
                    v = parse_data_blob(r['data'])
                    if isinstance(v, dict):
                        v = v.get('USD') or next(iter(v.values()))
                    ser[r['d'].strftime('%Y-%m-%d')] = float(v)
                except Exception:
                    continue
            out[key] = ser
    finally:
        conn.close()
    return out


def gen_chart(data, end):
    W, H = 1420, 720
    PL, PT_, PW, PH = 90, 24, 1250, 560
    common_d = sorted(set(data['gold']) & set(data['growth']) & set(data['fx']))
    gold = [(d, data['gold'][d] * data['fx'][d]) for d in common_d]
    grw = [(d, data['growth'][d] * data['fx'][d]) for d in common_d]
    g0, w0 = gold[0][1], grw[0][1]
    gl = [(d, math.log(v / g0) * 100) for d, v in gold]
    wl = [(d, math.log(v / w0) * 100) for d, v in grw]
    allv = [v for _, v in gl] + [v for _, v in wl]
    YMAX = math.ceil(max(allv) / 50) * 50
    YMIN = math.floor(min(allv) / 50) * 50
    T0, T1 = _dnum(common_d[0]), _dnum(common_d[-1])

    def X(d): return PL + (_dnum(d) - T0) / (T1 - T0) * PW
    def Y(v): return PT_ + (YMAX - v) / (YMAX - YMIN) * PH

    fig = plt.figure(figsize=(W / 72, H / 72), dpi=144)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W); ax.set_ylim(H, 0); ax.axis('off')
    ax.add_patch(plt.Rectangle((PL, PT_), PW, PH, facecolor='none',
                               edgecolor='#BFBFBF', lw=1.2, zorder=1))
    FS = 19                                       # 눈금·범례 = 10pt (2026-07-13 사용자 지시)
    for v in range(YMIN, YMAX + 1, 50):
        ax.plot([PL, PL + PW], [Y(v)] * 2, color='#9a9a9a' if v == 0 else '#EBEBEB',
                lw=1, zorder=1)
        ax.text(PL - 10, Y(v), f'{v}%', ha='right', va='center', fontsize=FS, color='#333')
    y0, y1 = int(common_d[0][:4]), int(common_d[-1][:4])
    for yr in range(y0 - 1, y1 + 1):
        d = f'{yr}-12-31'
        if _dnum(d) < T0 or _dnum(d) > T1:
            continue
        ax.plot([X(d)] * 2, [PT_ + PH, PT_ + PH + 7], color='#BFBFBF', lw=1)
        ax.text(X(d) + 6, PT_ + PH + 18, d, fontsize=FS, color='#555',
                rotation=32, ha='right', va='top', rotation_mode='anchor')
    ax.plot([X(d) for d, _ in gl], [Y(v) for _, v in gl], color=C_GOLD, lw=1.8, zorder=3)
    ax.plot([X(d) for d, _ in wl], [Y(v) for _, v in wl], color=C_GROWTH, lw=1.8, zorder=3)
    # 이벤트 타원·라벨 위치만 산출 — 도형/텍스트는 pptx 네이티브(편집 가능, 2026-07-13 지시)
    gmap, wmap = dict(gl), dict(wl)
    events_geo = []
    for ed, label in EVENTS:
        d = next((x for x in common_d if x >= ed), None)
        if d is None:
            continue
        ya, yb = Y(gmap[d]), Y(wmap[d])
        cy, hh = (ya + yb) / 2, abs(ya - yb) + 70
        events_geo.append({'cx': X(d), 'cy': cy, 'w': 46, 'h': hh,
                           'label': label, 'label_y': min(ya, yb) - 46})
    # 범례 (하단 중앙, 10pt)
    ly = PT_ + PH + 96
    lx = W / 2 - 200
    ax.plot([lx, lx + 26], [ly] * 2, color=C_GOLD, lw=3.4)
    ax.text(lx + 33, ly, '금(KRW)', fontsize=FS, color='#333', ha='left', va='center')
    lx += 33 + 150 + 26
    ax.plot([lx, lx + 26], [ly] * 2, color=C_GROWTH, lw=3.4)
    ax.text(lx + 33, ly, '미국성장주(KRW)', fontsize=FS, color='#333', ha='left', va='center')
    fig.savefig(OUT / 's13_chart.png', facecolor='white')
    plt.close(fig)
    return OUT / 's13_chart.png', (W, H), events_geo


def add(prs, ctx, page_label='13'):
    end = ctx['asof']
    sub = ('미국성장주와 금은 기대수익률은 높으면서도 낮은 상관관계로 분산투자효과가 높음. '
           '(두 자산이 모두 하락하는 빈도 낮음)')
    sl = slide_scaffold(prs, 'base_slide07.png', '미국 성장주와 금의 분산투자효과', end,
                        page_label, subtitle=sub)
    add_text(sl, 100, 208, 1400, 34, '누적로그수익률: 금 vs 미국성장주', 24 * PT_PER_PX,
             '222222', bold=True, align=PP_ALIGN.CENTER)
    data = _load(end)
    png, (w, h), events_geo = gen_chart(data, end)
    x0, y0 = (1600 - w) // 2, 258
    sl.shapes.add_picture(str(png), EX(x0), E(y0), E(w), E(h))
    # 이벤트 타원(점선) + 라벨 = 네이티브 편집 개체 (2026-07-13 사용자 지시)
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    for ev in events_geo:
        sh = sl.shapes.add_shape(
            MSO_SHAPE.OVAL,
            EX(x0 + ev['cx'] - ev['w'] / 2), E(y0 + ev['cy'] - ev['h'] / 2),
            E(ev['w']), E(ev['h']))
        sh.fill.background()
        sh.shadow.inherit = False
        sh.line.color.rgb = RGBColor.from_string('555555')
        sh.line.width = Pt(1.0)
        ln = sh.line._get_or_add_ln()
        pd = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(pd)
        add_text(sl, x0 + ev['cx'] - 120, y0 + ev['label_y'] - 26, 240, 26,
                 ev['label'], BODY_PT - 2, 'B8860B', align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.BOTTOM)
    add_text(sl, 900, 1015, 660, 26, f'· 자료: {kdate(end)}, Factset, 한국투자신탁운용',
             BODY_PT - 2, '787878', align=PP_ALIGN.RIGHT)
    return sl
