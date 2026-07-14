"""운용보고 PPT 빌더 — 공통 상수·A4 셸·pptx 헬퍼.

P0' 파일럿(pilot/pilot_s7.py)에서 확정한 스펙을 패키지화:
  - A4 가로 29.7x21cm, base 1600x900 → 1600x1131 (y500 흰행 삽입, 헤더/푸터 무왜곡)
  - 콘텐츠 y196~1055 선형 재배치 SV=1.3678, 본문 12pt 고정
  - 한글 a:ea/a:cs 슬롯, 표 셀 tcPr@anchor, 빈 셀 문단 폰트, ko-KR
  - 밴드 전폭 균일화 + 코너딥 평탄화, 로고/제목/페이지번호 = 편집 개체로 분리
"""
import sys
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
from matplotlib import font_manager
import matplotlib.pyplot as plt

BUILDER = Path(__file__).parent
ROOT = BUILDER.parent                     # reporting/
REPO = ROOT.parent                        # DB_OCIO_Webview
if str(REPO) not in sys.path:
    sys.path.insert(0, str(REPO))         # modules/, config/ import 용

OUT = ROOT / 'out'
FONTS = ROOT / 'template' / 'fonts'
KI_LOGO = REPO / 'web' / 'public' / 'ki-logo.png'

for _f in FONTS.glob('*.otf'):
    font_manager.fontManager.addfont(str(_f))
plt.rcParams['font.family'] = 'Pretendard'

# ── 페이지 좌표계 (2026-07-14 사용자 재구성 템플릿 780x540pt 기준) ──
# 내부 설계 px(1600x1131, 구 A4 좌표)는 전부 유지하고 출력만 균등 축소:
# 1px = 540pt/1131 = 0.47745pt (높이 맞춤) + 가로 중앙정렬 X_OFF(+8.05pt).
# 사용자 파일 역산 검증: 제목 36.27pt=76px, 본문 10.89pt, 이미지/표 위치 일치.
EMU_PAGE_W, EMU_PAGE_H = 780 * 12_700, 540 * 12_700  # 780 x 540 pt
EMU_PER_PX = EMU_PAGE_H / 1131                       # 6063.66
X_OFF = EMU_PAGE_W / EMU_PER_PX / 2 - 800            # 16.86px = 8.05pt (가로 중앙)
PX_H = 1131
PT_PER_PX = 540 / 1131                               # 0.47745 (px→pt)
BODY_PT = round(12 * PT_PER_PX / 0.5262, 2)          # 10.89 (구 A4 12pt 등가)
CANVAS_OFF = (40, 192)
_FOOT_OLD, _FOOT_NEW, _HDR = 824, 824 + (PX_H - 900), 196   # 푸터라인 1055
SV = (_FOOT_NEW - _HDR) / (_FOOT_OLD - _HDR)                 # 1.3678

# 색 (발송본 팔레트)
HDR_BLUE = '5B9BD5'; Z1 = 'D2DEEF'; Z2 = 'EAEFF7'
C_FUND = '#2E5E9E'; C_BM = '#E0A800'; C_EXC = '#B9C9E8'
BROWN = '7B401F'; INK = '222222'; RED = 'C00000'; BLUE = '0563C1'
PHASE_NAME = {1: '회복', 2: '팽창', 3: '침체', 4: '둔화'}


def remap(y):
    """900px 레이아웃 y → A4(1131px) 콘텐츠존 y."""
    return round(_HDR + (y - _HDR) * SV) if y >= _HDR else y


def sv(h):
    """세로 길이 스케일."""
    return round(h * SV)


# ──────────────────────────── pptx 헬퍼 ────────────────────────────
from pptx import Presentation                                  # noqa: E402
from pptx.util import Emu, Pt                                  # noqa: E402
from pptx.dml.color import RGBColor                            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                         # noqa: E402
from pptx.enum.lang import MSO_LANGUAGE_ID                     # noqa: E402
from pptx.oxml.ns import qn                                    # noqa: E402
from lxml import etree                                         # noqa: E402


def E(px):
    """길이(폭·높이) px → EMU."""
    return Emu(round(px * EMU_PER_PX))


def EX(px):
    """x 위치 px → EMU (가로 중앙정렬 오프셋 포함)."""
    return Emu(round((px + X_OFF) * EMU_PER_PX))


def set_ko_font(font, family):
    """a:latin 만으로는 한글이 테마 EA(맑은고딕)로 폴백 — a:ea/a:cs 도 지정.

    멱등: 같은 rPr 에 재호출해도 a:ea/a:cs 가 중복 삽입되지 않음
    (중복되면 스키마 위반으로 PowerPoint 가 파일을 거부 — 2026-07-14 확인).
    """
    font.language_id = MSO_LANGUAGE_ID.KOREAN
    rPr = font._rPr
    latin = rPr.find(qn('a:latin')) if rPr is not None else None
    if latin is None:
        return
    for tag in ('a:cs', 'a:ea'):            # addnext 역순 → latin, ea, cs 순서
        old = rPr.find(qn(tag))
        if old is not None:
            rPr.remove(old)
        e = rPr.makeelement(qn(tag), {'typeface': family})
        latin.addnext(e)


def PTE(pt):
    """pt → EMU (정적 슬라이드·레이아웃 — 사용자 파일 실측 pt 좌표 직접 사용)."""
    return Emu(round(pt * 12_700))


def new_presentation():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(EMU_PAGE_W), Emu(EMU_PAGE_H)
    _setup_layouts(prs)
    return prs


def layout_by_name(prs, name):
    return prs.slide_masters[0].slide_layouts.get_by_name(name)


_FRAME_DISCLAIMER = (
    '본 자료는 당사의 승인 없이 불법적으로 복제 또는 유통될 수 없습니다. '
    '본 자료에 기재된 운용 전략 및 전망은 시장상황 변동 등에 따라 변경될 수 '
    '있으며, 상기 예시 수익률이 미래의 수익을 보장하는 것은 아닙니다. '
    '본 자료 중 예측 및 전망에 관한 자료는 참고 자료이며 향후의 결과를 '
    '보증하는 것은 아닙니다.')

_NS = ('xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
       'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"')


def _frame_shape_xmls():
    """'1_Title Slide' 레이아웃 프레임(백색 캔버스+푸터라인+면책문구) sp XML —
    사용자 재구성 템플릿 실측(pt) 그대로 (2026-07-14)."""
    def _e(pt):
        return round(pt * 12_700)

    rect = (
        f'<p:sp {_NS}><p:nvSpPr><p:cNvPr id="901" name="Frame Canvas"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>'
        f'<a:xfrm><a:off x="{_e(0)}" y="{_e(86)}"/>'
        f'<a:ext cx="{_e(780)}" cy="{_e(454)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:p/></p:txBody></p:sp>')
    line = (
        f'<p:cxnSp {_NS}><p:nvCxnSpPr><p:cNvPr id="902" name="Frame Footer Line"/>'
        f'<p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr><p:spPr>'
        f'<a:xfrm><a:off x="{_e(22.8)}" y="{_e(504.7)}"/>'
        f'<a:ext cx="{_e(734.4)}" cy="0"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="9525"><a:solidFill><a:srgbClr val="D9D9D9"/></a:solidFill></a:ln>'
        f'</p:spPr></p:cxnSp>')
    disc = (
        f'<p:sp {_NS}><p:nvSpPr><p:cNvPr id="903" name="Frame Disclaimer"/>'
        f'<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr>'
        f'<a:xfrm><a:off x="{_e(25.2)}" y="{_e(508.8)}"/>'
        f'<a:ext cx="{_e(686.5)}" cy="{_e(30.5)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0"/>'
        f'<a:p><a:r><a:rPr lang="ko-KR" sz="800" dirty="0">'
        f'<a:solidFill><a:srgbClr val="4A2119"/></a:solidFill>'
        f'<a:latin typeface="Pretendard"/><a:ea typeface="Pretendard"/>'
        f'<a:cs typeface="Pretendard"/></a:rPr>'
        f'<a:t>{_FRAME_DISCLAIMER}</a:t></a:r></a:p></p:txBody></p:sp>')
    return [rect, line, disc]


def _add_layout(prs, name, shape_xmls):
    """slideLayout 파트를 새로 만들어 마스터에 등록 (python-pptx 에 API 없음 —
    blank 레이아웃 XML 을 복제해 이름/도형 교체, placeholder 는 제거)."""
    from copy import deepcopy
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.packuri import PackURI
    from pptx.oxml import parse_xml
    from pptx.parts.slide import SlideLayoutPart

    master = prs.slide_masters[0]
    blank = prs.slide_layouts[6]
    el = deepcopy(blank._element)
    cSld = el.find(qn('p:cSld'))
    cSld.set('name', name)
    spTree = cSld.find(qn('p:spTree'))
    for sp in spTree.findall(qn('p:sp')):        # 기본 placeholder 제거 (사용자 레이아웃 동일)
        spTree.remove(sp)
    for xml in shape_xmls:
        spTree.append(parse_xml(xml))

    existing = {str(ly.part.partname) for m in prs.slide_masters for ly in m.slide_layouts}
    idx = 12
    while f'/ppt/slideLayouts/slideLayout{idx}.xml' in existing:
        idx += 1
    part = SlideLayoutPart(PackURI(f'/ppt/slideLayouts/slideLayout{idx}.xml'),
                           blank.part.content_type, prs.part.package, el)
    part.relate_to(master.part, RT.SLIDE_MASTER)
    rid = master.part.relate_to(part, RT.SLIDE_LAYOUT)
    lst = master._element.find(qn('p:sldLayoutIdLst'))
    new_id = max(int(e.get('id')) for e in lst) + 1
    ent = lst.makeelement(qn('p:sldLayoutId'), {'id': str(new_id)})
    ent.set(qn('r:id'), rid)
    lst.append(ent)


def _setup_layouts(prs):
    """사용자 재구성 템플릿(2026-07-14)의 커스텀 레이아웃 3종 생성.

    '1_Title Slide'   — 데이터 슬라이드 프레임 (백색 캔버스·푸터라인·면책문구)
    '1_Title and Content' — 섹션 표지 (우측 백색 패널)
    '제목 및 내용'      — 표지·목차 (빈 레이아웃)
    """
    def _e(pt):
        return round(pt * 12_700)

    right_panel = (
        f'<p:sp {_NS}><p:nvSpPr><p:cNvPr id="911" name="Section Right Panel"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>'
        f'<a:xfrm><a:off x="{_e(264.2)}" y="0"/>'
        f'<a:ext cx="{_e(515.8)}" cy="{_e(540)}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:p/></p:txBody></p:sp>')
    _add_layout(prs, '1_Title Slide', _frame_shape_xmls())
    _add_layout(prs, '1_Title and Content', [right_panel])
    _add_layout(prs, '제목 및 내용', [])


def add_text(sl, px_x, px_y, px_w, px_h, text, pt_size, color, bold=False,
             family='Pretendard', align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=False):
    tb = sl.shapes.add_textbox(EX(px_x), E(px_y), E(px_w), E(px_h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = family; r.font.size = Pt(pt_size)
    r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
    set_ko_font(r.font, family)
    return tb


def add_bullets(sl, px_x, px_y, px_w, lines, pt_size=BODY_PT, color=INK,
                line_gap_px=None):
    """여러 불릿 문단 텍스트박스 (word_wrap, 문단 간격)."""
    from pptx.util import Pt as _Pt
    tb = sl.shapes.add_textbox(EX(px_x), E(px_y), E(px_w), E(40 * len(lines)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = _Pt((line_gap_px or sv(14)) * PT_PER_PX / 0.75)  # px→pt 근사
        r = p.add_run(); r.text = line
        r.font.name = 'Pretendard'; r.font.size = Pt(pt_size)
        r.font.color.rgb = RGBColor.from_string(color)
        set_ko_font(r.font, 'Pretendard')
    return tb


def add_pbar(sl, px_x, px_y, px_w, text, px_h=None):
    """파란 소제목 바 (편집 가능 도형)."""
    h = px_h or sv(34)
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, EX(px_x), E(px_y), E(px_w), E(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(HDR_BLUE)
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = 'Pretendard'; r.font.size = Pt(BODY_PT)
    r.font.bold = True; r.font.color.rgb = RGBColor.from_string('FFFFFF')
    set_ko_font(r.font, 'Pretendard')
    return sh


def add_table(sl, px_x, px_y, col_w_px, row_h_px, rows_spec):
    """네이티브 표. rows_spec: [(fill_hex|None, [(text, bold, color_hex), ...]), ...]"""
    n_r, n_c = len(rows_spec), len(col_w_px)
    gf = sl.shapes.add_table(n_r, n_c, EX(px_x), E(px_y),
                             E(sum(col_w_px)), E(sum(row_h_px)))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is not None:
        sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'   # No Style, No Grid
    for i, w in enumerate(col_w_px):
        tbl.columns[i].width = E(w)
    for i, h in enumerate(row_h_px):
        tbl.rows[i].height = E(h)
    for ri, (fill, cells) in enumerate(rows_spec):
        for ci, (text, bold, color) in enumerate(cells):
            cell = tbl.cell(ri, ci)
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(fill)
            if ri < n_r - 1:                          # 행 사이 2px 흰 경계
                tcPr = cell._tc.get_or_add_tcPr()
                ln = etree.SubElement(tcPr, qn('a:lnB'))
                ln.set('w', '15240'); ln.set('cap', 'flat')
                sf = etree.SubElement(ln, qn('a:solidFill'))
                etree.SubElement(sf, qn('a:srgbClr')).set('val', 'FFFFFF')
                tcPr.insert(0, ln)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE   # 표 셀 세로가운데 = tcPr@anchor
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Emu(round(2 * EMU_PER_PX))
            tf.margin_top = tf.margin_bottom = 0
            tf.word_wrap = False
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Pretendard'; p.font.size = Pt(BODY_PT)   # 빈 셀 행높이 방지
            set_ko_font(p.font, 'Pretendard')
            if text:
                r = p.add_run(); r.text = text
                r.font.name = 'Pretendard'; r.font.size = Pt(BODY_PT)
                r.font.bold = bold; r.font.color.rgb = RGBColor.from_string(color)
                set_ko_font(r.font, 'Pretendard')
    return gf


def kdate(iso):
    y, m, d = iso.split('-')
    return f'{y}년 {int(m)}월 {int(d)}일'


def slide_scaffold(prs, base_name, title, asof_iso, page_label, subtitle='기준일'):
    """공통 골격: 고해상 로고 + 제목 + 부제 + 페이지번호.

    base_name: (구 A4 셸 PNG — 2026-07-14 사용자 템플릿 재구성으로 폐지, 호환용 무시.
    프레임(백색 캔버스·푸터라인·면책문구)은 '1_Title Slide' 레이아웃에 있음.)
    subtitle: '기준일'(기본 — asof 로 포맷) | 임의 문자열 | None(생략, s11/12 등)
    """
    sl = prs.slides.add_slide(layout_by_name(prs, '1_Title Slide'))
    LOGO_W = 150
    LOGO_H = round(LOGO_W * 607 / 2467)
    sl.shapes.add_picture(str(KI_LOGO), EX(1600 - 27 - LOGO_W), E(14), E(LOGO_W), E(LOGO_H))
    add_text(sl, 36, 22, 1000, 96, title, 76 * PT_PER_PX, '000000', family='Pretendard Black')
    if subtitle == '기준일':
        subtitle = f'기준일: {kdate(asof_iso)}'
    if subtitle:
        add_text(sl, 55, 133, 1100, 34, subtitle, 26 * PT_PER_PX, BROWN, bold=True)
    add_text(sl, 1440, 1062, 120, 30, page_label, 24 * PT_PER_PX, '444444',
             align=PP_ALIGN.RIGHT)
    return sl
