#!/usr/bin/env python3
"""202606 운용보고 PPT 데이터 슬라이드 빌더.

데이터 스냅샷이 갱신되면 이 스크립트 하나로 4·6·7·9·10·11·12p를 재생성한다.

    cd report_202606/build && ../../../pptx_env/bin/python build_slides.py

입력:
  - ../../valuation.js            (PER/EPS/price 일별 시계열, window.VALUATION)
  - ../data/fund_07G07_YTD.json   (펀드 NAV/BM/보유)
  - ../data/fund_07G07_2Q.json
  - slide04_data.json             (4p 자산군표·코멘트 — 수동 관리)
  - base/base_slideNN.png         (데이터 영역이 비어있는 고정 템플릿)

출력:
  - ../valuation_report_20260630.pptx 의 해당 슬라이드 blob 교체
  - ../ppt_report_20260630/slideNN.png 캡처 동기화

렌더링: playwright(chromium) — market-brief/node_modules 사용.
"""
import json, re, subprocess, datetime, sys
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

HERE = Path(__file__).parent
ROOT = HERE.parent            # report_202606
VALROOT = ROOT.parent         # Valuation
OUT = HERE / 'out'
FONTS = ROOT / 'slide_src' / 'fonts'
NODE_DIR = VALROOT / 'market-brief'   # playwright가 설치된 node_modules 위치
PPTX = ROOT / 'valuation_report_20260630.pptx'
CAPT = ROOT / 'ppt_report_20260630'

T0, T1 = None, None  # 차트 기간 (아래에서 설정)
def dnum(d): y,m,dd = map(int, d.split('-')); return datetime.date(y,m,dd).toordinal()

FF = '''  @font-face { font-family:'Pretendard'; src:url('FONTDIR/Pretendard-Regular.otf'); font-weight:400; }
  @font-face { font-family:'Pretendard'; src:url('FONTDIR/Pretendard-Medium.otf'); font-weight:500; }
  @font-face { font-family:'Pretendard'; src:url('FONTDIR/Pretendard-SemiBold.otf'); font-weight:600; }
  @font-face { font-family:'Pretendard'; src:url('FONTDIR/Pretendard-Bold.otf'); font-weight:700; }
'''.replace('FONTDIR', str(FONTS))
BROWN = (123, 64, 31)
SRC_FS = '· 자료: {}, Factset, 한국투자신탁운용'
SRC_BB = '· 자료: {}, Bloomberg, 한국투자신탁운용'

def font(w, s): return ImageFont.truetype(str(FONTS / f'Pretendard-{w}.otf'), s)

# ──────────────────────────── 데이터 로드 ────────────────────────────
def load_all():
    s = (VALROOT / 'valuation.js').read_text()
    val = {e['sym']: e for e in json.loads(s[len('window.VALUATION='):].rstrip().rstrip(';'))}
    fund = json.loads((ROOT / 'data' / 'fund_07G07_YTD.json').read_text())
    fund2q = json.loads((ROOT / 'data' / 'fund_07G07_2Q.json').read_text())
    s4 = json.loads((HERE / 'slide04_data.json').read_text())
    return val, fund, fund2q, s4

def kdate(iso):
    y, m, d = iso.split('-')
    return f'{y}년 {int(m)}월 {int(d)}일'

# ──────────────────────────── s4: 자산군표 + 코멘트 ────────────────────────────
def gen_s4(s4):
    mx = s4['bar_max']
    def cell_bar(v, klass, override=None):
        if v is None:
            return f'<td class="{klass} dash">&ndash;</td>'
        label = override or f'{v:.2f}%'
        return f'<td class="bar {klass}" data-v="{v}" data-l="{label}"></td>'
    rows_html = []
    for r in s4['rows']:
        tds = []
        if r.get('cat'):
            tds.append(f'<td class="cat" rowspan="{r["cat_span"]}">{r["cat"]}</td>')
        if r['type'] == 'grp':
            tds.append(f'<td colspan="2">{r["label"]}</td>')
        else:
            if r.get('sub'):
                cls = 'sub' if r['type'] == 'row' else 'sub'
                tds.append(f'<td class="{cls}" rowspan="{r["sub_span"]}">{r["sub"]}</td>')
            k = 'd' if r['type'] == 'krow' else ''
            if r.get('label_colspan'):
                tds.append(f'<td class="{k}" colspan="{r["label_colspan"]}">{r["label"]}</td>')
            else:
                tds.append(f'<td class="{k}">{r["label"]}</td>')
        k = 'd' if r['type'] == 'krow' else ''
        tds.append(cell_bar(r['v1'], k, r.get('label_override')))
        tds.append(cell_bar(r['v2'], k, r.get('label_override')))
        tds.append(f'<td class="bench {k}">{r["bench"]}</td>')
        cls = {'grp': ' class="grp"', 'krow': ' class="krow"', 'row': ''}[r['type']]
        rows_html.append(f'  <tr{cls}>{"".join(tds)}</tr>')
    table = f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ width:858px; background:#fff; font-family:'Pretendard',sans-serif; }}
  table {{ width:858px; border-collapse:collapse; table-layout:fixed; border:1px solid #B7B7B7; }}
  td, th {{ border:1px solid #B7B7B7; font-size:12.5px; color:#000; text-align:center; vertical-align:middle; height:18.4px; padding:0 2px; overflow:hidden; white-space:nowrap; }}
  th {{ background:#5B9BD5; color:#fff; font-weight:700; font-size:12.5px; line-height:1.15; }}
  tr.hdr th {{ height:33px; }}
  .cat {{ font-weight:700; font-size:13px; background:#fff; }}
  .grp td {{ background:#F2F2F2; font-weight:700; }}
  tr.krow td.d, tr.krow td.sub {{ background:#F8E6D6; }}
  td.sub {{ background:#fff; }}
  .bench {{ font-size:12px; }}
  .dash {{ color:#7f7f7f; font-weight:400; }}
  td.bar {{ position:relative; padding:0; }}
  td.bar .bg {{ position:absolute; left:1px; top:2px; bottom:2px;
    background:linear-gradient(to right,#6CC68B 0%,#A8DDBC 60%,#E9F6EE 100%); border:1px solid #4EA36A; border-left:none; }}
  td.bar .bg.neg {{ background:linear-gradient(to right,#FF5860 0%,#FFB3B7 70%,#FFE9EA 100%); border-color:#D64550; }}
  td.bar .tx {{ position:relative; z-index:1; font-size:12.5px; }}
  td.bar::before {{ content:''; position:absolute; left:0; top:0; bottom:0; border-left:1px dashed #AAA; }}
  .grp td.bar .tx {{ font-weight:700; }}
</style>
<table>
  <colgroup>
    <col style="width:83px"><col style="width:118px"><col style="width:161px">
    <col style="width:97px"><col style="width:97px"><col style="width:302px">
  </colgroup>
  <tr class="hdr">
    <th colspan="3">자산군</th>
    <th>기간 수익률<br>(표시 통화)</th>
    <th>기간 수익률<br>(원화환산)</th>
    <th>벤치마크</th>
  </tr>
{chr(10).join(rows_html)}
</table>
<script>
  const MAX = {mx};
  document.querySelectorAll('td.bar').forEach(td => {{
    const v = parseFloat(td.dataset.v);
    const w = Math.max(Math.abs(v) / MAX * (td.clientWidth - 3), Math.abs(v) > 0.01 ? 3 : 0);
    const bg = document.createElement('div');
    bg.className = 'bg' + (td.dataset.v.startsWith('-') ? ' neg' : '');
    bg.style.width = w + 'px';
    if (w === 0) bg.style.display = 'none';
    const tx = document.createElement('div');
    tx.className = 'tx'; tx.textContent = td.dataset.l;
    td.appendChild(bg); td.appendChild(tx);
  }});
</script>'''
    secs = []
    for i, sec in enumerate(s4['comments']):
        lines = '\n      '.join(f'<div>{l}</div>' for l in sec['lines'])
        secs.append(f'''  <div class="sec" style="min-height:{sec["min_h"]}px{';border-top:1px solid #E4E4E4' if i==0 else ''}">
    <div class="lbl">{sec["label"]}</div>
    <div class="txt">
      {lines}
    </div>
  </div>''')
    comment = f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ width:670px; background:#fff; font-family:'Pretendard',sans-serif; }}
  #wrap {{ width:670px; }}
  .sec {{ display:flex; border-bottom:1px solid #E4E4E4; align-items:center; }}
  .lbl {{ width:112px; flex:none; text-align:center; font-size:19px; font-weight:700; color:#111; line-height:1.35; }}
  .txt {{ flex:1; font-size:19px; font-weight:500; color:#1a1a1a; line-height:27px; padding:10px 0 10px 2px; word-break:keep-all; }}
  .txt div {{ padding-left:15px; text-indent:-15px; }}
</style>
<div id="wrap">
{chr(10).join(secs)}
</div>'''
    return table, comment

# ──────────────────────────── s6/s7: 펀드 ────────────────────────────
BM_WEIGHTS = {'국내주식': 0.0, '해외주식': 34.0, '국내채권': 41.0, '해외채권': 25.0, '현금': 0.0}
# 근거: DB_OCIO/config/funds.py FUND_BM['07G04'] (07G07=07G04 클래스펀드)
#       0.34×MSCI ACWI Gross + 0.25×BBG AGG(H) + 0.41×KIS KTB10Y

def fund_calcs(fund, fund2q):
    sm, sm2 = fund['summary'], fund2q['summary']
    ser = fund['series']
    end = ser[-1]
    end_d = datetime.date(*map(int, end['date'].split('-')))
    py, pm = (end_d.year, end_d.month - 1) if end_d.month > 1 else (end_d.year - 1, 12)
    last_dom = (datetime.date(py, pm, 28) + datetime.timedelta(days=4)).replace(day=1) - datetime.timedelta(days=1)
    m1_target = datetime.date(py, pm, min(end_d.day, last_dom.day))
    m1 = next(r for r in reversed(ser) if r['date'] <= m1_target.isoformat())
    weights = {}
    for x in fund['securities']:
        weights[x['bucket']] = weights.get(x['bucket'], 0) + x['weight']
    cash = 100 - sum(weights.values())
    return {
        'asof': sm['end'], 'asof_kr': kdate(sm['end']),
        'ytd_f': sm['fund_return_pct'], 'ytd_b': sm['bm_return_pct'], 'ytd_x': sm['excess_pct'],
        'q_f': sm2['fund_return_pct'], 'q_b': sm2['bm_return_pct'], 'q_x': sm2['excess_pct'],
        'm1_f': round((end['nav']/m1['nav']-1)*100, 2), 'm1_b': round((end['bm']/m1['bm']-1)*100, 2),
        'si_f': round((end['nav']/1000-1)*100, 2), 'si_b': round((end['bm']/1000-1)*100, 2),
        'w': {k: round(v, 1) for k, v in weights.items()}, 'cash': round(cash, 1),
        'series': ser,
    }

TBL_CSS = '''  table { position:absolute; border-collapse:separate; border-spacing:0 2px; table-layout:fixed; }
  th { background:#5B9BD5; color:#fff; font-weight:700; font-size:18px; vertical-align:middle; }
  td { font-size:18px; color:#222; text-align:center; vertical-align:middle; }
  tr.z1 td { background:#D2DEEF; }
  tr.z2 td { background:#EAEFF7; }
'''

def gen_s6(fc):
    w = fc['w']; bw = BM_WEIGHTS
    def arow(cls, name, f, b):
        act = f - b
        return f'<tr class="{cls}"><td>{name}</td><td>{f:.1f}</td><td>{b:.1f}</td><td>{act:+.1f}</td></tr>'
    rows = [
        arow('z1', '국내주식', w.get('국내주식', 0), bw['국내주식']),
        arow('z2', '해외주식', w.get('해외주식', 0), bw['해외주식']),
        arow('z1', '국내채권', w.get('국내채권', 0), bw['국내채권']),
        arow('z2', '해외채권', w.get('해외채권', 0), bw['해외채권']),
        arow('z1', '현금', fc['cash'], bw['현금']),
    ]
    return f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ font-family:'Pretendard',sans-serif; background:#fff; }}
  #c {{ position:relative; width:1520px; height:620px; }}
  .bl {{ position:absolute; left:14px; font-size:19px; color:#222; }}
{TBL_CSS}
  th {{ height:54px; }} td {{ height:56px; }}
  td.b {{ font-weight:700; }}
</style>
<div id="c">
  <div class="bl" style="top:22px">· 연초 이후 수익률: 펀드 {fc['ytd_f']:+.2f}%, BM {fc['ytd_b']:+.2f}% (초과성과 {fc['ytd_x']:+.2f}%p)</div>
  <div class="bl" style="top:54px">· 2분기 수익률: 펀드 {fc['q_f']:+.2f}%, BM {fc['q_b']:+.2f}% (초과성과 {fc['q_x']:+.2f}%p)</div>
  <div class="bl" style="top:86px">· 보유자산 배분: 국내채권 {fc['w'].get('국내채권',0):.1f}%(국고채 10·30년 중심), 해외주식 {fc['w'].get('해외주식',0):.1f}%, 국내주식 {fc['w'].get('국내주식',0):.1f}% (BM: KTB10Y 41 / ACWI 34 / 글로벌채권(H) 25)</div>
  <table style="left:24px; top:192px; width:592px">
    <tr><th>기준 월</th><th>관측 월</th><th>경기국면</th></tr>
    <tr class="z1"><td>2024년 8월</td><td>2024년 9월</td><td>둔화</td></tr>
    <tr class="z2"><td>2024년 10월</td><td>2024년 11월</td><td>팽창</td></tr>
    <tr class="z1"><td>2025년 2월</td><td>2025년 3월</td><td>둔화</td></tr>
    <tr class="z2"><td class="b">2025년 6월</td><td class="b">2025년 7월</td><td class="b">팽창</td></tr>
    <tr class="z1"><td class="b">2026년 3월</td><td class="b">2026년 4월</td><td class="b">둔화</td></tr>
  </table>
  <table style="left:656px; top:192px; width:754px">
    <colgroup><col style="width:150px"><col style="width:190px"><col style="width:190px"><col style="width:224px"></colgroup>
    <tr><th>자산군</th><th>펀드 비중(%)</th><th>BM 비중(%)</th><th>Active 비중(%p)</th></tr>
    {chr(10).join(rows)}
  </table>
</div>'''

def gen_s7(fc):
    ser = fc['series']
    n0, b0 = ser[0]['nav'], ser[0]['bm']
    CW, CH = 640, 280
    ymax, ymin = 12, -2
    def Y(v): return (ymax - v)/(ymax - ymin)*CH
    N = len(ser)
    pts_f, pts_b, pts_e, mt = [], [], [], []
    for i, r in enumerate(ser):
        x = i/(N-1)*CW
        f = (r['nav']/n0-1)*100; b = (r['bm']/b0-1)*100
        pts_f.append(f'{x:.1f},{Y(f):.1f}'); pts_b.append(f'{x:.1f},{Y(b):.1f}')
        pts_e.append((x, f-b))
        if r['date'][8:10] == '01':
            mt.append((x, str(int(r['date'][5:7])) + '월'))
    area = f'M0,{Y(0):.1f} ' + ' '.join(f'L{x:.1f},{Y(e):.1f}' for x, e in pts_e) + f' L{CW},{Y(0):.1f} Z'
    gl = ''.join(f'<line x1="0" y1="{Y(v):.0f}" x2="{CW}" y2="{Y(v):.0f}" stroke="{"#999" if v==0 else "#E8E8E8"}" stroke-width="1"/>' for v in range(ymin, ymax+1, 2))
    ylab = ''.join(f'<div style="position:absolute; left:-40px; top:{Y(v)-9:.0f}px; width:34px; text-align:right; font-size:13px; color:#666">{v}%</div>' for v in range(ymin, ymax+1, 2))
    xlab = ''.join(f'<div style="position:absolute; left:{x-14:.0f}px; top:{CH+6}px; font-size:13px; color:#666">{m}</div>' for x, m in mt)
    si_f = f'{fc["si_f"]:.2f}' if fc.get('si_f') is not None else ''
    si_b = f'{fc["si_b"]:.2f}' if fc.get('si_b') is not None else ''
    return f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ font-family:'Pretendard',sans-serif; background:#fff; }}
  #c {{ position:relative; width:1520px; height:625px; }}
  .bl {{ position:absolute; left:14px; font-size:19px; color:#222; }}
{TBL_CSS}
  th {{ height:38px; }} td {{ height:36px; }}
  td.rl {{ font-weight:700; }}
  .pbar {{ position:absolute; height:34px; background:#5B9BD5; color:#fff; font-weight:700; font-size:18px; text-align:center; line-height:34px; }}
  .legend {{ font-size:13.5px; color:#555; }}
</style>
<div id="c">
  <div class="bl" style="top:22px">· 연초 이후 수익률 {fc['ytd_f']:+.2f}%로 BM({fc['ytd_b']:+.2f}%) 대비 {fc['ytd_x']:+.2f}%p 초과 성과</div>
  <div class="bl" style="top:54px">· 2분기 수익률: 펀드 {fc['q_f']:+.2f}%, BM {fc['q_b']:+.2f}% (초과성과 {fc['q_x']:+.2f}%p)</div>
  <table style="left:52px; top:96px; width:1430px">
    <colgroup><col style="width:230px"><col style="width:240px"><col style="width:240px"><col style="width:240px"><col style="width:240px"><col style="width:240px"></colgroup>
    <tr><th>기간</th><th>1개월</th><th>3개월</th><th>6개월</th><th>연초 이후</th><th>설정 이후</th></tr>
    <tr class="z1"><td class="rl">펀드(%)</td><td>{fc['m1_f']:.2f}</td><td>{fc['q_f']:.2f}</td><td>{fc['ytd_f']:.2f}</td><td>{fc['ytd_f']:.2f}</td><td>{si_f}</td></tr>
    <tr class="z2"><td class="rl">BM(%)</td><td>{fc['m1_b']:.2f}</td><td>{fc['q_b']:.2f}</td><td>{fc['ytd_b']:.2f}</td><td>{fc['ytd_b']:.2f}</td><td>{si_b}</td></tr>
  </table>
  <div class="pbar" style="left:48px; top:246px; width:724px">연초 이후 성과 추이</div>
  <div style="position:absolute; left:96px; top:300px; width:{CW}px; height:{CH}px;">
    <svg width="{CW}" height="{CH}" style="position:absolute; overflow:visible">
      {gl}
      <path d="{area}" fill="#B9C9E8" fill-opacity="0.8"/>
      <polyline points="{' '.join(pts_f)}" fill="none" stroke="#2E5E9E" stroke-width="2"/>
      <polyline points="{' '.join(pts_b)}" fill="none" stroke="#E0A800" stroke-width="2"/>
    </svg>
    {ylab}{xlab}
    <div class="legend" style="position:absolute; left:110px; top:{CH+28}px;">
      <span style="display:inline-block;width:22px;height:11px;background:#B9C9E8;margin-right:5px;vertical-align:-1px"></span>초과성과(%p)
      <span style="display:inline-block;width:22px;height:3px;background:#2E5E9E;margin:0 5px 3px 22px"></span>펀드
      <span style="display:inline-block;width:22px;height:3px;background:#E0A800;margin:0 5px 3px 22px"></span>BM
    </div>
  </div>
  <div class="pbar" style="left:788px; top:246px; width:694px">기여수익률 분석(단위: %, %p)</div>
  <table style="left:788px; top:288px; width:694px">
    <colgroup><col style="width:190px"><col style="width:168px"><col style="width:168px"><col style="width:168px"></colgroup>
    <tr><th>연초 이후</th><th>펀드</th><th>BM</th><th>Active</th></tr>
    <tr class="z1"><td>국내주식</td><td></td><td></td><td></td></tr>
    <tr class="z2"><td>해외주식</td><td></td><td></td><td></td></tr>
    <tr class="z1"><td>국내채권</td><td></td><td></td><td></td></tr>
    <tr class="z2"><td>해외채권</td><td></td><td></td><td></td></tr>
    <tr class="z1"><td>대체투자</td><td></td><td></td><td></td></tr>
    <tr class="z2"><td>기타</td><td></td><td></td><td></td></tr>
    <tr class="z1"><td class="rl">합계</td><td class="rl">{fc['ytd_f']:+.2f}</td><td class="rl">{fc['ytd_b']:+.2f}</td><td class="rl">{fc['ytd_x']:+.2f}</td></tr>
  </table>
</div>'''

# ──────────────────────────── s9~s12: 밸류에이션 차트 ────────────────────────────
S9_CATS = [('전세계', 'MXWD Index'), ('미국', 'MXUS Index'), ('미국<br>빅테크 7 Plus', None),
           ('미국<br>성장주', 'MXUS000G'), ('미국<br>가치주', 'MXUS000V'), ('DM ex US', 'MXWOU Index'),
           ('EM', 'MXEF Index'), ('한국', 'MXKR Index'), ('중국', 'MXCN Index'), ('일본', 'MXJP Index'),
           ('독일', None), ('영국', 'MXGB Index'), ('브라질', None), ('남아공', None)]
S9_SUBTITLE_SCOPE = [('DM ex US', 'MXWOU Index'), ('US Growth', 'MXUS000G'), ('US', 'MXUS Index'),
                     ('Global', 'MXWD Index'), ('US Value', 'MXUS000V'), ('EM', 'MXEF Index'), ('Korea', 'MXKR Index')]

def s9_decomp(val, start, end):
    out = {}
    for nm, sym in dict(S9_CATS | dict(S9_SUBTITLE_SCOPE).items() if False else {}).items():
        pass
    syms = {sym for _, sym in S9_CATS if sym} | {sym for _, sym in S9_SUBTITLE_SCOPE}
    for sym in syms:
        ser = val[sym]['series']
        s0 = next(x for x in ser if x[0] >= start)
        s1 = next(x for x in reversed(ser) if x[0] <= end)
        perc = (s1[1]/s0[1]-1)*100
        epsc = (s1[2]/s0[2]-1)*100
        tot = ((1+perc/100)*(1+epsc/100)-1)*100
        out[sym] = (round(perc, 1), round(epsc, 1), round(tot-perc-epsc, 1), round(tot, 1))
    return out

def gen_s9(val, start, end):
    dec = s9_decomp(val, start, end)
    entries = []
    for nm, sym in S9_CATS:
        if sym:
            p, e, o, t = dec[sym]
            entries.append(f"  ['{nm}', {p}, {e}, {o}],")
        else:
            entries.append(f"  ['{nm}', null, null, null],")
    title = f'글로벌 주식 총 수익률 분석(USD, {start.replace("-","")} ~ {end.replace("-","")})'
    return f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ font-family:'Pretendard',sans-serif; background:#fff; }}
  #chart {{ position:relative; width:1520px; height:632px; }}
  .title {{ position:absolute; top:0; left:0; width:100%; text-align:center; font-size:24px; font-weight:700; color:#222; }}
  .gl {{ position:absolute; height:1px; background:#DDD; }}
  .gl.zero {{ background:#7f7f7f; height:2px; }}
  .yt {{ position:absolute; width:66px; text-align:right; font-size:16px; color:#333; transform:translateY(-50%); left:36px; }}
  .seg {{ position:absolute; }}
  table {{ position:absolute; border-collapse:collapse; table-layout:fixed; }}
  td, th {{ border:1px solid #C9C9C9; font-size:14.5px; color:#222; text-align:center; vertical-align:middle; padding:2px 1px; line-height:1.15; overflow:hidden; }}
  th {{ font-weight:700; height:50px; background:#F5F5F5; }}
  td.rl {{ font-weight:700; text-align:left; padding-left:8px; background:#FAFAFA; font-size:15px; }}
  tr.tot td {{ font-weight:700; background:#F2F2F2; }}
  td {{ height:41px; font-size:15px; }}
  .sw {{ display:inline-block; width:14px; height:14px; border-radius:2px; margin-right:6px; vertical-align:-2px; }}
  .fn {{ position:absolute; font-size:13px; color:#777; left:112px; }}
</style>
<div id="chart">
  <div class="title">{title}</div>
</div>
<script>
const C = {{ B:'#2C5A9C', R:'#A83232', G:'#6E8B3D' }};
const DATA = [
{chr(10).join(entries)}
];
const val = c => c;
const chart = document.getElementById('chart');
const ML = 112, MR = 17, PT = 40, PB = 378;
const PW = 1520 - ML - MR, PH = PB - PT;
const YMAX = 250, YMIN = -100, S = PH / (YMAX - YMIN);
const y0 = PT + YMAX * S;
for (let v = 250; v >= -100; v -= 50) {{
  const y = PT + (YMAX - v) * S;
  const gl = document.createElement('div');
  gl.className = 'gl' + (v === 0 ? ' zero' : '');
  gl.style.cssText = `left:${{ML}}px; width:${{PW}}px; top:${{v===0?y-1:y}}px`;
  chart.appendChild(gl);
  const t = document.createElement('div'); t.className = 'yt';
  t.style.top = y + 'px'; t.textContent = v + '%';
  chart.appendChild(t);
}}
const slot = PW / DATA.length, BW = 56;
DATA.forEach((row, i) => {{
  const cx = ML + slot * (i + 0.5);
  let up = 0, dn = 0;
  [['B',row[1]],['R',row[2]],['G',row[3]]].forEach(([k, v]) => {{
    if (v === null) return;
    const h = Math.abs(v) * S;
    let top;
    if (v >= 0) {{ top = y0 - (up + v) * S; up += v; }}
    else {{ top = y0 + dn * S; dn += -v; }}
    const d = document.createElement('div');
    d.className = 'seg';
    d.style.cssText = `left:${{cx - BW/2}}px; top:${{top}}px; width:${{BW}}px; height:${{h}}px; background:${{C[k]}}`;
    chart.appendChild(d);
  }});
}});
const fmt = v => v === null ? '' : v.toFixed(1) + '%';
const totCell = r => (r[1]===null && r[2]===null && r[3]===null) ? '' : (r[1]+r[2]+r[3]).toFixed(1) + '%';
const tbl = document.createElement('table');
tbl.style.cssText = `left:0px; top:${{PB + 8}}px; width:${{112 + (1520-112-17)}}px;`;
let html = '<colgroup><col style="width:112px">' +
  DATA.map(() => `<col style="width:${{(1520-112-17)/DATA.length}}px">`).join('') + '</colgroup>';
html += '<tr><th style="background:#fff;border-top:none;border-left:none"></th>' + DATA.map(r => `<th>${{r[0]}}</th>`).join('') + '</tr>';
html += '<tr class="tot"><td class="rl">총수익</td>' + DATA.map(r => `<td>${{totCell(r)}}</td>`).join('') + '</tr>';
html += `<tr><td class="rl"><span class="sw" style="background:${{C.B}}"></span>PER 변화율</td>` + DATA.map(r => `<td>${{fmt(r[1])}}</td>`).join('') + '</tr>';
html += `<tr><td class="rl"><span class="sw" style="background:${{C.R}}"></span>EPS 변화율</td>` + DATA.map(r => `<td>${{fmt(r[2])}}</td>`).join('') + '</tr>';
html += `<tr><td class="rl"><span class="sw" style="background:${{C.G}}"></span>기타</td>` + DATA.map(r => `<td>${{fmt(r[3])}}</td>`).join('') + '</tr>';
tbl.innerHTML = html;
chart.appendChild(tbl);
const fn = document.createElement('div');
fn.className = 'fn';
fn.style.top = (PB + 8 + 50 + 41*4 + 12) + 'px';
fn.textContent = '총수익은 12M 선행 PER×EPS 기반 가격수익률(USD, 배당 제외), 기타는 PER·EPS 변동의 교차항 · 빈칸은 데이터 입수 시 반영 예정';
chart.appendChild(fn);
const src = document.createElement('div');
src.style.cssText = 'position:absolute; right:17px; top:' + (PB + 8 + 50 + 41*4 + 12) + 'px; font-size:15px; color:#787878;';
src.textContent = 'SRCLINE';
chart.appendChild(src);
</script>'''.replace('SRCLINE', SRC_BB.format(kdate(end)))

def s9_subtitle(val, start, end):
    dec = s9_decomp(val, start, end)
    ranked = sorted(S9_SUBTITLE_SCOPE, key=lambda t: dec[t[1]][3])
    parts = [nm for nm, _ in ranked]
    return ' < '.join(parts[:-1]) + ' <<< ' + parts[-1]

def gen_s10(val, end):
    PL, PT, PW, PH = 78, 44, 1372, 480
    kr = val['MXKR Index']
    def X(d): return PL + (dnum(d)-T0)/(T1-T0)*PW
    def YL(v): return PT + (4400-v)/4000*PH
    def YR(p): return PT + (22-p)/16*PH
    red = ' '.join(f'{X(d):.1f},{YL(v):.1f}' for d, v in kr['price'])
    blue = ' '.join(f'{X(d):.1f},{YR(p):.1f}' for d, p, e in kr['series'])
    d_b, p_b = next((d, p) for d, p, e in kr['series'] if d >= '2021-03-02')
    d_r, v_r = next((d, v) for d, v in kr['price'] if d >= '2019-09-02')
    yl = ''.join(f'<div class="yt l" style="top:{PT+PH*i/8:.0f}px">{4400-i*500:,}</div>' for i in range(9))
    yl += ''.join(f'<div class="yt r" style="top:{PT+PH*i/8:.0f}px">{(22-i*2):.1f}</div>' for i in range(9))
    xt = ''.join(f'<div class="tick" style="left:{X(f"{2015+i}-09-30"):.0f}px"></div><div class="xt" style="left:{X(f"{2015+i}-09-30")+6:.0f}px">{2015+i}-09-30</div>' for i in range(11))
    grid = ''.join(f'<line x1="{PL}" y1="{PT+PH*i/8:.0f}" x2="{PL+PW}" y2="{PT+PH*i/8:.0f}" stroke="#EBEBEB" stroke-width="1"/>' for i in range(1, 8))
    return f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ font-family:'Pretendard',sans-serif; background:#fff; }}
  #chart {{ position:relative; width:1520px; height:640px; }}
  .title {{ position:absolute; top:0; left:{PL}px; width:{PW}px; text-align:center; font-size:24px; font-weight:700; color:#222; }}
  .yt {{ position:absolute; font-size:19px; color:#333; transform:translateY(-50%); }}
  .yt.l {{ width:66px; text-align:right; left:4px; }}
  .yt.r {{ left:1458px; }}
  .xt {{ position:absolute; top:{PT+PH+22}px; font-size:18px; color:#333;
    transform:rotate(-32deg) translateX(-100%); transform-origin:left top; white-space:nowrap; }}
  .tick {{ position:absolute; top:{PT+PH}px; width:1.2px; height:7px; background:#BFBFBF; }}
  .ann {{ position:absolute; font-size:19px; line-height:1.25; text-align:center; padding:5px 12px; background:#fff; white-space:nowrap; z-index:3; }}
  .ann.b {{ border:2px solid #2E5E9E; color:#2E5E9E; }}
  .ann.r {{ border:2px solid #A83232; color:#A83232; }}
  .src {{ position:absolute; right:70px; top:618px; font-size:15px; color:#787878; }}
</style>
<div id="chart">
  <div class="title">한국 가격지수 vs 밸류에이션</div>
  <svg width="1520" height="640">
    <rect x="{PL}" y="{PT}" width="{PW}" height="{PH}" fill="none" stroke="#BFBFBF" stroke-width="1.2"/>
    {grid}
    <polyline points="{red}" fill="none" stroke="#C0392B" stroke-width="2"/>
    <polyline points="{blue}" fill="none" stroke="#2E6DB4" stroke-width="2"/>
    <line x1="960" y1="112" x2="{X(d_b):.0f}" y2="{YR(p_b)-4:.0f}" stroke="#2E5E9E" stroke-width="2" stroke-dasharray="7,5"/>
    <line x1="420" y1="256" x2="{X(d_r):.0f}" y2="{YL(v_r)-4:.0f}" stroke="#A83232" stroke-width="2" stroke-dasharray="7,5"/>
  </svg>
  <div class="ann b" style="left:910px; top:56px;">12개월 선행 PER<br>( MSCI Korea, 우)</div>
  <div class="ann r" style="left:255px; top:198px;">Price Index<br>( MSCI Korea, 좌)</div>
  {yl}{xt}
  <div class="src">{SRC_BB.format(kdate(end))}</div>
</div>'''

S11_SERIES = [('MXKR Index', 'Korea', '#A6A6A6', '#7F7F7F'), ('MXUS000G', 'US Growth', '#FFC000', '#D99E00'),
              ('MXUS Index', 'US', '#ED7D31', '#E06B1F'), ('MXUS000V', 'US Value', '#4472C4', '#3A66B8'),
              ('MXEF Index', 'EM', '#5B9BD5', '#4E95D4'), ('MXWOU Index', 'DM ex US', '#1F3864', '#1F3864')]

def gen_s11(val, end):
    PL2, PT2, PW2, PH2 = 84, 44, 1230, 470
    def X2(d): return PL2 + (dnum(d)-T0)/(T1-T0)*PW2
    def Y2(v): return PT2 + (800-v)/900*PH2
    lines, ends = [], []
    for sym, nm, col, tcol in S11_SERIES:
        rows = val[sym]['series']; e0 = rows[0][2]
        pts = ' '.join(f'{X2(d):.1f},{Y2((e/e0-1)*100):.1f}' for d, p, e in rows)
        lines.append(f'<polyline points="{pts}" fill="none" stroke="{col}" stroke-width="2"/>')
        endv = (rows[-1][2]/e0-1)*100
        ends.append((nm, endv, Y2(endv), tcol))
    labh, prev = [], -1e9
    for nm, v, y, tcol in sorted(ends, key=lambda t: t[2]):
        y = max(y, prev + 25); prev = y
        labh.append(f'<div class="slab" style="top:{y:.0f}px;color:{tcol}">{nm} {v:.0f}%</div>')
    yl2 = ''.join(f'<div class="yt" style="top:{PT2+PH2*i/9:.0f}px">{800-i*100}%</div>' for i in range(10))
    grid2 = ''.join(f'<line x1="{PL2}" y1="{PT2+PH2*i/9:.0f}" x2="{PL2+PW2}" y2="{PT2+PH2*i/9:.0f}" stroke="{"#9a9a9a" if 800-i*100==0 else "#EBEBEB"}" stroke-width="1"/>' for i in range(10))
    xt2 = ''.join(f'<div class="tick" style="left:{X2(f"{2015+i}-09-30"):.0f}px"></div><div class="xt" style="left:{X2(f"{2015+i}-09-30")+6:.0f}px">{2015+i}-09-30</div>' for i in range(11))
    return f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ font-family:'Pretendard',sans-serif; background:#fff; }}
  #chart {{ position:relative; width:1520px; height:640px; }}
  .title {{ position:absolute; top:0; left:84px; width:1230px; text-align:center; font-size:24px; font-weight:700; color:#222; }}
  .yt {{ position:absolute; width:72px; text-align:right; left:6px; font-size:19px; color:#333; transform:translateY(-50%); }}
  .xt {{ position:absolute; top:{PT2+PH2+18}px; font-size:18px; color:#333;
    transform:rotate(-32deg) translateX(-100%); transform-origin:left top; white-space:nowrap; }}
  .tick {{ position:absolute; top:{PT2+PH2}px; width:1.2px; height:7px; background:#BFBFBF; }}
  .slab {{ position:absolute; left:1322px; font-size:19px; font-weight:700; white-space:nowrap; transform:translateY(-50%); }}
  .src {{ position:absolute; right:{1520-(PL2+PW2)}px; top:616px; font-size:15px; color:#787878; }}
</style>
<div id="chart">
  <div class="title">12개월 선행 EPS 추이</div>
  <svg width="1520" height="640">{grid2}{''.join(lines)}</svg>
  {yl2}{xt2}{''.join(labh)}
  <div class="src">{SRC_FS.format(kdate(end))}</div>
</div>'''

S12_CATS = [('전세계', 'MXWD Index'), ('미국', 'MXUS Index'), ('미국<br>빅테크<br>7 Plus', None),
            ('미국<br>성장주', 'MXUS000G'), ('미국<br>가치주', 'MXUS000V'), ('DM ex<br>US', 'MXWOU Index'),
            ('EM', 'MXEF Index'), ('한국', 'MXKR Index'), ('중국', 'MXCN Index'), ('일본', 'MXJP Index'),
            ('독일', None), ('영국', 'MXGB Index'), ('인도', None), ('러시아', None), ('브라질', None),
            ('인니', None), ('남아공', None)]

def gen_s12(val, start, end):
    PL3, PT3, PW3, PH3 = 60, 66, 1400, 440
    def Y12(v): return PT3 + (52-v)/52*PH3
    marks, xlabs = [], []
    slot = PW3/len(S12_CATS)
    for i, (nm, sym) in enumerate(S12_CATS):
        cx = PL3 + slot*(i+0.5)
        xlabs.append(f'<div class="xt" style="left:{cx:.0f}px">{nm}</div>')
        if not sym: continue
        pers = [r[1] for r in val[sym]['series']]
        mn, mx, me, cu = min(pers), max(pers), sum(pers)/len(pers), pers[-1]
        marks.append(f'<rect x="{cx-13:.0f}" y="{Y12(mx):.1f}" width="26" height="{Y12(mn)-Y12(mx):.1f}" rx="10" fill="#D9DEB0"/>')
        marks.append(f'<line x1="{cx-19:.0f}" y1="{Y12(me):.1f}" x2="{cx+19:.0f}" y2="{Y12(me):.1f}" stroke="#111" stroke-width="3.2"/>')
        marks.append(f'<circle cx="{cx:.0f}" cy="{Y12(mx):.1f}" r="7.5" fill="#D9412B"/>')
        marks.append(f'<circle cx="{cx:.0f}" cy="{Y12(mn):.1f}" r="7.5" fill="#3D9BDC"/>')
        marks.append(f'<rect x="{cx-7:.0f}" y="{Y12(cu)-7:.1f}" width="14" height="14" fill="#1F3352" transform="rotate(45 {cx:.0f} {Y12(cu):.1f})"/>')
    grid3 = ''.join(f'<line x1="{PL3}" y1="{Y12(v):.0f}" x2="{PL3+PW3}" y2="{Y12(v):.0f}" stroke="#EBEBEB" stroke-width="1"/>' for v in range(0, 51, 10))
    ylabs3 = ''.join(f'<div class="yt" style="top:{Y12(v):.0f}px">{v}</div>' for v in range(0, 51, 10))
    title = f'지역/국가별 12개월 선행 PER({start.replace("-","")} ~ {end.replace("-","")})'
    return f'''<meta charset="utf-8">
<style>
{FF}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ font-family:'Pretendard',sans-serif; background:#fff; }}
  #chart {{ position:relative; width:1520px; height:640px; }}
  .title {{ position:absolute; top:8px; left:60px; width:1400px; text-align:center; font-size:23px; font-weight:700; color:#222; }}
  .yt {{ position:absolute; width:44px; text-align:right; left:8px; font-size:19px; color:#333; transform:translateY(-50%); }}
  .xt {{ position:absolute; text-align:center; font-size:16.5px; color:#222; line-height:1.15; transform:translateX(-50%); top:{PT3+PH3+8}px; }}
  .legend {{ position:absolute; top:576px; left:0; width:100%; text-align:center; font-size:18px; color:#222; }}
  .legend .it {{ margin:0 18px; }}
  .dot {{ display:inline-block; width:15px; height:15px; border-radius:50%; margin-right:7px; vertical-align:-2px; }}
  .bar {{ display:inline-block; width:26px; height:3.5px; background:#111; margin-right:7px; vertical-align:4px; }}
  .dia {{ display:inline-block; width:13px; height:13px; background:#1F3352; transform:rotate(45deg); margin-right:9px; vertical-align:-1px; }}
  .src {{ position:absolute; right:60px; top:616px; font-size:15px; color:#787878; }}
</style>
<div id="chart">
  <div class="title">{title}</div>
  <svg width="1520" height="640">
    <rect x="{PL3}" y="{PT3}" width="{PW3}" height="{PH3}" fill="none" stroke="#BFBFBF" stroke-width="1.2"/>
    {grid3}{''.join(marks)}
  </svg>
  {ylabs3}{''.join(xlabs)}
  <div class="legend">
    <span class="it"><span class="dot" style="background:#D9412B"></span>최대</span>
    <span class="it"><span class="dot" style="background:#3D9BDC"></span>최소</span>
    <span class="it"><span class="bar"></span>평균</span>
    <span class="it"><span class="dia"></span>현재</span>
  </div>
  <div class="src">{SRC_FS.format(kdate(end))}</div>
</div>'''

# ──────────────────────────── 렌더 & 합성 ────────────────────────────
def render(targets):
    shot = HERE / 'shot.js'
    shot.write_text('''const { chromium } = require(process.argv[2] + '/node_modules/playwright');
(async () => {
  const jobs = JSON.parse(process.argv[3]);
  const browser = await chromium.launch();
  const page = await browser.newPage({ viewport: { width: 1700, height: 1300 }, deviceScaleFactor: 2 });
  for (const [html, sel, out] of jobs) {
    await page.goto('file://' + html);
    await page.waitForTimeout(250);
    await (await page.$(sel)).screenshot({ path: out });
  }
  await browser.close();
})();
''')
    jobs = [[str(OUT / f'{k}.html'), sel, str(OUT / f'{k}.png')] for k, sel in targets]
    subprocess.run(['node', str(shot), str(NODE_DIR), json.dumps(jobs)], check=True, cwd=str(HERE))

def half(p):
    im = Image.open(p).convert('RGB')
    return im.resize((im.width//2, im.height//2), Image.LANCZOS)

def compose(fc, s9sub, end):
    finals = {}
    def base(n): return Image.open(HERE / 'base' / f'base_slide{n:02d}.png').convert('RGB')
    FB26 = font('Bold', 26)
    s4cfg = json.loads((HERE / 'slide04_data.json').read_text())
    # 4p
    img = base(4); d = ImageDraw.Draw(img)
    d.text((55, 133), s4cfg['period_line'], font=FB26, fill=BROWN)
    img.paste(half(OUT / 's4_table.png'), (38, 250))
    img.paste(half(OUT / 's4_comment.png'), (908, 250))
    finals[4] = img
    # 6p, 7p
    for n in (6, 7):
        img = base(n); d = ImageDraw.Draw(img)
        d.text((55, 133), f'기준일: {fc["asof_kr"]}', font=FB26, fill=BROWN)
        img.paste(half(OUT / f's{n}.png'), (40, 192))
        finals[n] = img
    # 9p~12p
    subs = {9: s9sub, 10: f'{end.replace("-",".").replace(".0",".")} 기준 MXKR PER(12FWD): {fc["mxkr_pe"]:.2f}'}
    for n in (9, 10, 11, 12):
        img = base(n); d = ImageDraw.Draw(img)
        if n in subs:
            d.text((55, 133), subs[n], font=FB26, fill=BROWN)
        r = half(OUT / f's{n}.png')
        if r.width > 1540 or r.height > 620:
            f = min(1540/r.width, 620/r.height)
            r = r.resize((int(r.width*f), int(r.height*f)), Image.LANCZOS)
        img.paste(r, ((1600 - r.width)//2, 200))
        finals[n] = img
    return finals

def write_pptx(finals):
    from pptx import Presentation
    prs = Presentation(str(PPTX))
    for n, img in finals.items():
        p = OUT / f'final_slide{n:02d}.png'
        img.save(p)
        sl = prs.slides[n-1]
        pic = list(sl.shapes)[0]
        part = sl.part.related_part(pic._element.blipFill.blip.rEmbed)
        part._blob = p.read_bytes()
        img.save(CAPT / f'slide{n:02d}.png')
    prs.save(str(PPTX))

def main():
    global T0, T1
    val, fund, fund2q, s4 = load_all()
    fc = fund_calcs(fund, fund2q)
    start, end = fund['summary']['start'], fund['summary']['end']
    kr_ser = val['MXKR Index']['series']
    T0, T1 = dnum(kr_ser[0][0]), dnum(kr_ser[-1][0])
    fc['mxkr_pe'] = kr_ser[-1][1]
    OUT.mkdir(exist_ok=True)
    t, c = gen_s4(s4)
    (OUT / 's4_table.html').write_text(t)
    (OUT / 's4_comment.html').write_text(c)
    (OUT / 's6.html').write_text(gen_s6(fc))
    (OUT / 's7.html').write_text(gen_s7(fc))
    (OUT / 's9.html').write_text(gen_s9(val, start, end))
    (OUT / 's10.html').write_text(gen_s10(val, end))
    (OUT / 's11.html').write_text(gen_s11(val, end))
    (OUT / 's12.html').write_text(gen_s12(val, kr_ser[0][0], end))
    print('HTML 생성 완료 → 렌더링...')
    render([('s4_table', 'table'), ('s4_comment', '#wrap'), ('s6', '#c'), ('s7', '#c'),
            ('s9', '#chart'), ('s10', '#chart'), ('s11', '#chart'), ('s12', '#chart')])
    print('렌더링 완료 → 합성/저장...')
    s9sub = s9_subtitle(val, start, end)
    print('  9p 부제:', s9sub)
    finals = compose(fc, s9sub, end)
    write_pptx(finals)
    print(f'완료: {PPTX.name} 슬라이드 {sorted(finals)} 갱신, 캡처 동기화')

if __name__ == '__main__':
    main()
